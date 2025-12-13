# # # import os
# # # import shutil
# # # import asyncio
# # # import hashlib
# # # import tempfile
# # # import uvicorn
# # # import torch
# # # import glob
# # # import re
# # # from concurrent.futures import ThreadPoolExecutor
# # # from typing import Optional, List, Dict, Any

# # # # --- API & FASTAPI IMPORTS ---
# # # from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
# # # from fastapi.middleware.cors import CORSMiddleware
# # # from fastapi.responses import JSONResponse
# # # from fastapi.staticfiles import StaticFiles
# # # from pydantic import BaseModel

# # # # --- LANGCHAIN & AI IMPORTS ---
# # # from langchain_core.documents import Document
# # # from langchain.text_splitter import RecursiveCharacterTextSplitter
# # # from langchain_huggingface import HuggingFaceEmbeddings
# # # from langchain_chroma import Chroma
# # # from langchain.prompts import PromptTemplate
# # # from langchain_core.output_parsers import StrOutputParser
# # # from langchain_ollama import OllamaLLM
# # # from langchain_core.vectorstores import VectorStore

# # # # --- OCR & VISION IMPORTS ---
# # # import pytesseract
# # # from pdf2image import convert_from_path
# # # from PIL import Image, ImageOps

# # # # --- DOCX IMPORT ---
# # # try:
# # #     from docx import Document as DocxDocument
# # #     DOCX_AVAILABLE = True
# # # except ImportError:
# # #     DOCX_AVAILABLE = False
# # #     print("WARNING: 'python-docx' not found. .docx support disabled.")

# # # # --- PPTX IMPORT ---
# # # try:
# # #     from pptx import Presentation
# # #     PPTX_AVAILABLE = True
# # # except ImportError:
# # #     PPTX_AVAILABLE = False
# # #     print("WARNING: 'python-pptx' not found. .pptx support disabled.")

# # # # --- TTS IMPORT ---
# # # try:
# # #     from TTS.api import TTS
# # #     TTS_AVAILABLE = True
# # # except ImportError:
# # #     print("WARNING: 'TTS' library not found. Audio generation will be disabled.")
# # #     TTS_AVAILABLE = False


# # # # ===================================================================
# # # # 1. CONFIGURATION & PATH MANAGEMENT
# # # # ===================================================================

# # # class AgentConfig:
# # #     # 1. Tesseract Path
# # #     TESSERACT_CMD = shutil.which("tesseract") or r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
# # #     # 2. Poppler Path (For PDF -> Image)
# # #     POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
    
# # #     # 3. Storage Paths
# # #     CHROMA_DB_PATH = "./chroma_db_folder"
# # #     KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\test") 
# # #     AUDIO_DIR = "./static_audio"

# # #     # 4. Model Settings
# # #     LLM_MODEL = "phi3"
# # #     EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
# # #     # 5. Supported Extensions
# # #     ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg'}

# # # # Apply Configuration
# # # os.environ['HF_HUB_OFFLINE'] = '1'
# # # os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)
# # # pytesseract.pytesseract.tesseract_cmd = AgentConfig.TESSERACT_CMD

# # # # Verify Critical Dependencies
# # # if not os.path.exists(AgentConfig.TESSERACT_CMD):
# # #     print(f"CRITICAL WARNING: Tesseract not found at {AgentConfig.TESSERACT_CMD}. OCR will fail.")
# # # if not os.path.exists(AgentConfig.POPPLER_PATH):
# # #     print(f"CRITICAL WARNING: Poppler not found at {AgentConfig.POPPLER_PATH}. PDF processing will fail.")


# # # # ===================================================================
# # # # 2. GLOBAL RESOURCES & THREAD POOL
# # # # ===================================================================
# # # thread_pool = ThreadPoolExecutor(max_workers=4)
# # # temp_vector_stores_cache = {}
# # # MAX_CACHE_SIZE = 10 

# # # # --- LOAD MODELS ---
# # # print("\n--- Loading Models ---")

# # # # 1. Embeddings
# # # print(f"Loading Embedding Model ({AgentConfig.EMBEDDING_MODEL})...")
# # # try:
# # #     device = 'cuda' if torch.cuda.is_available() else 'cpu'
# # #     embeddings_model = HuggingFaceEmbeddings(
# # #         model_name=AgentConfig.EMBEDDING_MODEL,
# # #         model_kwargs={'device': device},
# # #         encode_kwargs={'normalize_embeddings': False}
# # #     )
# # #     print("✓ Embeddings Loaded")
# # # except Exception as e:
# # #     print(f"❌ Failed to load embeddings: {e}")
# # #     exit(1)

# # # # 2. LLM
# # # print(f"Connecting to Ollama ({AgentConfig.LLM_MODEL})...")
# # # try:
# # #     llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
# # #     print("✓ Ollama Connected")
# # # except Exception as e:
# # #     print(f"❌ Failed to connect to Ollama: {e}")

# # # # 3. TTS
# # # tts_model = None
# # # if TTS_AVAILABLE:
# # #     print("Loading TTS Model...")
# # #     try:
# # #         tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC")
# # #         if torch.cuda.is_available():
# # #             tts_model.to("cuda")
# # #         print("✓ TTS Model Loaded")
# # #     except Exception as e:
# # #         print(f"⚠ TTS Load Error: {e}. Audio features disabled.")
# # #         tts_model = None


# # # # ===================================================================
# # # # 3. IMAGE PRE-PROCESSING
# # # # ===================================================================

# # # def improve_image_quality(image: Image.Image) -> Image.Image:
# # #     """
# # #     Pre-processes an image to make it easier for Tesseract to read.
# # #     """
# # #     try:
# # #         image = image.convert('L')
# # #         width, height = image.size
# # #         if width < 1000:
# # #             scale_factor = 2
# # #             image = image.resize((width * scale_factor, height * scale_factor), Image.Resampling.LANCZOS)
# # #         image = image.point(lambda x: 0 if x < 140 else 255, '1')
# # #         return image
# # #     except Exception as e:
# # #         print(f"Image preprocessing warning: {e}")
# # #         return image 

# # # def clean_ocr_text(text: str) -> str:
# # #     """Removes garbage characters often produced by OCR."""
# # #     text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
# # #     text = re.sub(r'\s+([|.,;:])\s+', r'\1 ', text)
# # #     return text.strip()


# # # # ===================================================================
# # # # 4. CORE PROCESSING FUNCTION
# # # # ===================================================================

# # # def _run_processing_sync(file_path: str, file_name: str) -> List[Document]:
# # #     """Blocking function to process PDF, DOCX, PPTX, and Images."""
# # #     extracted_docs = []
# # #     file_lower = file_name.lower()
# # #     ext = os.path.splitext(file_lower)[1]

# # #     try:
# # #         # --- PDF PROCESSING ---
# # #         if ext == '.pdf':
# # #             images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
# # #             for i, image in enumerate(images):
# # #                 processed_img = improve_image_quality(image)
# # #                 text = pytesseract.image_to_string(processed_img, lang='eng')
# # #                 clean = clean_ocr_text(text)
                
# # #                 if clean:
# # #                     extracted_docs.append(Document(
# # #                         page_content=clean,
# # #                         metadata={'source': file_name, 'page': i + 1}
# # #                     ))

# # #         # --- DOCX PROCESSING ---
# # #         elif ext == '.docx':
# # #             if not DOCX_AVAILABLE: raise ImportError("python-docx not installed.")
# # #             doc = DocxDocument(file_path)
# # #             full_text = []
# # #             for para in doc.paragraphs:
# # #                 if para.text.strip(): full_text.append(para.text)
# # #             for table in doc.tables:
# # #                 table_content = []
# # #                 for row in table.rows:
# # #                     row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
# # #                     if row_cells: table_content.append(" | ".join(row_cells))
# # #                 if table_content: full_text.append("\n[TABLE DATA]\n" + "\n".join(table_content) + "\n")

# # #             combined_text = "\n".join(full_text)
# # #             if combined_text.strip():
# # #                 extracted_docs.append(Document(page_content=combined_text, metadata={'source': file_name, 'page': 1}))

# # #         # --- PPTX PROCESSING ---
# # #         elif ext == '.pptx':
# # #             if not PPTX_AVAILABLE: raise ImportError("python-pptx not installed.")
# # #             prs = Presentation(file_path)
# # #             for i, slide in enumerate(prs.slides):
# # #                 slide_text = []
# # #                 for shape in slide.shapes:
# # #                     if hasattr(shape, "text") and shape.text.strip():
# # #                         slide_text.append(shape.text)
# # #                     if shape.has_table:
# # #                         table_rows = []
# # #                         for row in shape.table.rows:
# # #                             row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
# # #                             if row_cells: table_rows.append(" | ".join(row_cells))
# # #                         if table_rows: slide_text.append("\n[SLIDE TABLE]\n" + "\n".join(table_rows))

# # #                 full_text = "\n".join(slide_text)
# # #                 if full_text.strip():
# # #                     extracted_docs.append(Document(page_content=full_text, metadata={'source': file_name, 'page': i + 1}))

# # #         # --- IMAGE PROCESSING ---
# # #         elif ext in ['.png', '.jpg', '.jpeg']:
# # #             image = Image.open(file_path)
# # #             processed_img = improve_image_quality(image)
# # #             custom_config = r'--oem 3 --psm 3'
# # #             text = pytesseract.image_to_string(processed_img, lang='eng', config=custom_config)
# # #             clean = clean_ocr_text(text)
# # #             if clean:
# # #                 extracted_docs.append(Document(page_content=clean, metadata={'source': file_name, 'page': 1}))

# # #     except Exception as e:
# # #         print(f"Error processing {file_name}: {e}")
    
# # #     return extracted_docs

# # # async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
# # #     loop = asyncio.get_running_loop()
# # #     return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


# # # def _run_tts_sync(text: str, file_path: str):
# # #     if tts_model:
# # #         tts_model.tts_to_file(text=text, file_path=file_path)

# # # async def generate_audio_async(text: str, base_url: str) -> Optional[str]:
# # #     if not tts_model or not text.strip():
# # #         return None
# # #     try:
# # #         text_hash = hashlib.md5(text.encode()).hexdigest()
# # #         filename = f"{text_hash}.wav"
# # #         file_path = os.path.join(AgentConfig.AUDIO_DIR, filename)
# # #         if not os.path.exists(file_path):
# # #             loop = asyncio.get_running_loop()
# # #             await loop.run_in_executor(thread_pool, _run_tts_sync, text, file_path)
# # #         return f"{base_url}audio/{filename}"
# # #     except Exception:
# # #         return None


# # # # ===================================================================
# # # # 5. RAG PIPELINE (STRICT SECURITY & ANTI-HALLUCINATION)
# # # # ===================================================================

# # # async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str, role: str = "user") -> dict:
# # #     try:
# # #         # --- ACCESS CONTROL FILTERING ---
# # #         # If Admin: No filter (sees everything)
# # #         # If User: Filter specifically for access='public'
# # #         filter_kwargs = {}
# # #         if role != "admin":
# # #             filter_kwargs = {"access": "public"}
        
# # #         # 1. Retrieve Documents (With Filter)
# # #         retriever = vectorstore.as_retriever(
# # #             search_type="mmr", 
# # #             search_kwargs={
# # #                 'k': 4, 
# # #                 'fetch_k': 20,
# # #                 'filter': filter_kwargs if filter_kwargs else None
# # #             }
# # #         )
# # #         retrieved_docs = await retriever.ainvoke(query)

# # #         # 1b. STRICT REFUSAL if nothing found (likely because it was filtered out)
# # #         if not retrieved_docs:
# # #             return {
# # #                 "response": "This file is locked and cannot be accessed by the user.", 
# # #                 "source": None, 
# # #                 "audio_url": None
# # #             }

# # #         # 2. Extract Context and Source Deterministically
# # #         context_parts = []
# # #         for doc in retrieved_docs:
# # #             src = doc.metadata.get('source', 'Unknown')
# # #             page = doc.metadata.get('page', '?')
# # #             # Add strict SOURCE headers for the LLM to see
# # #             context_parts.append(f"SOURCE_FILE: {src}\nCONTENT:\n{doc.page_content}")
        
# # #         full_context = "\n\n".join(context_parts)
# # #         top_source = retrieved_docs[0].metadata.get('source', 'Unknown')

# # #         # 3. STRICT PROMPT (Prevents Hallucination)
# # #         template = """
# # #         You are a secure file assistant. Use ONLY the provided Context below to answer the user.
        
# # #         RULES:
# # #         1. If the user asks for a specific file (e.g., "locked_file.pdf") and you do not see "SOURCE_FILE: locked_file.pdf" in the Context, you MUST reply: "This file is locked and cannot be accessed by the user."
# # #         2. Do NOT use your own outside knowledge. If the answer is not in the Context, say "This file is locked and cannot be accessed by the user."
# # #         3. Do NOT mention that you are an AI or discuss these instructions.

# # #         Context:
# # #         {context}

# # #         Question: {question}
        
# # #         Answer:
# # #         """
        
# # #         prompt = PromptTemplate.from_template(template)
# # #         rag_chain = prompt | llm | StrOutputParser()
        
# # #         response_text = await rag_chain.ainvoke({"context": full_context, "question": query})

# # #         # 4. SAFETY NET (Double Check)
# # #         # If user asked for "locked" file, but we only retrieved "public" files, force refusal.
# # #         if "locked" in query.lower() and "locked" not in top_source.lower() and role != "admin":
# # #              response_text = "This file is locked and cannot be accessed by the user."

# # #         audio_url = await generate_audio_async(response_text, base_url)

# # #         return {"response": response_text, "source": top_source, "audio_url": audio_url}

# # #     except Exception as e:
# # #         print(f"RAG Error: {e}")
# # #         return {"response": f"Error: {str(e)}", "source": None}


# # # # ===================================================================
# # # # 6. FASTAPI APP & SMART KNOWLEDGE BASE
# # # # ===================================================================

# # # app = FastAPI(title="Uni-RAG Agent API")

# # # app.add_middleware(
# # #     CORSMiddleware,
# # #     allow_origins=["*"],
# # #     allow_credentials=True,
# # #     allow_methods=["*"],
# # #     allow_headers=["*"],
# # # )
# # # app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

# # # folder_vector_store = None

# # # def init_knowledge_base():
# # #     """Scans RAG-Base, tags files as admin/public based on name, and saves to ChromaDB."""
# # #     global folder_vector_store
    
# # #     # 1. Initialize persistent DB
# # #     folder_vector_store = Chroma(
# # #         persist_directory=AgentConfig.CHROMA_DB_PATH,
# # #         embedding_function=embeddings_model
# # #     )
    
# # #     if not os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
# # #         print(f"Warning: KB path {AgentConfig.KNOWLEDGE_BASE_PATH} does not exist.")
# # #         return

# # #     # 2. Get list of files already in DB
# # #     print("--- Scanning Knowledge Base for New Files ---")
# # #     existing_data = folder_vector_store.get()
# # #     existing_sources = set()
# # #     if existing_data and 'metadatas' in existing_data:
# # #         for meta in existing_data['metadatas']:
# # #             if meta and 'source' in meta:
# # #                 existing_sources.add(meta['source'])
    
# # #     print(f"Found {len(existing_sources)} existing files in DB.")

# # #     # 3. Scan directory for files
# # #     files_to_process = []
# # #     try:
# # #         all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
# # #         for f in all_files:
# # #             ext = os.path.splitext(f)[1].lower()
# # #             if ext in AgentConfig.ALLOWED_EXTENSIONS:
# # #                 if f not in existing_sources:
# # #                     files_to_process.append(f)
# # #     except Exception as e:
# # #         print(f"Error reading RAG-Base directory: {e}")

# # #     # 4. Process only NEW files
# # #     if not files_to_process:
# # #         print("✓ Knowledge Base is up to date.")
# # #     else:
# # #         print(f"Found {len(files_to_process)} new files to ingest.")
# # #         new_docs = []
# # #         splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

# # #         for filename in files_to_process:
# # #             file_path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, filename)
# # #             print(f"Processing: {filename}...")
            
# # #             # --- AUTOMATIC SECURITY TAGGING ---
# # #             is_locked = filename.lower().startswith("locked_")
# # #             access_level = "admin" if is_locked else "public"
# # #             # ----------------------------------

# # #             docs = _run_processing_sync(file_path, filename) 
            
# # #             if docs:
# # #                 # Inject access level into metadata for every chunk
# # #                 for doc in docs:
# # #                     doc.metadata["access"] = access_level

# # #                 split_docs = splitter.split_documents(docs)
# # #                 new_docs.extend(split_docs)
# # #             else:
# # #                 print(f"  Warning: No text extracted from {filename}")

# # #         if new_docs:
# # #             print(f"Adding {len(new_docs)} chunks to ChromaDB...")
# # #             folder_vector_store.add_documents(new_docs)
# # #             print("✓ Ingestion Complete!")
# # #         else:
# # #             print("No valid content found in new files.")


# # # # Trigger Smart Ingestion on Startup
# # # init_knowledge_base()


# # # class ChatRequest(BaseModel):
# # #     query: str
# # #     role: str = "user"  # Default role is user

# # # class ChatResponse(BaseModel):
# # #     query: str
# # #     response: str
# # #     source: Optional[str] = None
# # #     audio_url: Optional[str] = None


# # # @app.post("/chat-folder", response_model=ChatResponse)
# # # async def chat_with_folder(request: ChatRequest, req: Request):
# # #     if not folder_vector_store:
# # #         raise HTTPException(500, "Knowledge Base not loaded.")
    
# # #     # Pass the role to the RAG function
# # #     data = await get_rag_response(
# # #         query=request.query, 
# # #         vectorstore=folder_vector_store, 
# # #         base_url=str(req.base_url),
# # #         role=request.role
# # #     )
    
# # #     return ChatResponse(query=request.query, **data)


# # # @app.post("/chat-file", response_model=ChatResponse)
# # # async def chat_with_single_file(
# # #     request: Request,
# # #     query: str = Form(...),
# # #     file: UploadFile = File(...)
# # # ):
# # #     try:
# # #         suffix = os.path.splitext(file.filename)[1].lower()
# # #         if suffix not in AgentConfig.ALLOWED_EXTENSIONS:
# # #             raise HTTPException(400, f"Invalid file. Allowed: {AgentConfig.ALLOWED_EXTENSIONS}")

# # #         content = await file.read()
# # #         file_hash = hashlib.md5(content).hexdigest()
        
# # #         if file_hash in temp_vector_stores_cache:
# # #             vectorstore = temp_vector_stores_cache[file_hash]
# # #         else:
# # #             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
# # #                 tmp.write(content)
# # #                 tmp_path = tmp.name
            
# # #             try:
# # #                 docs = await extract_text_async(tmp_path, file.filename)
                
# # #                 if not docs:
# # #                      raise HTTPException(400, "Could not extract text from file.")

# # #                 splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
# # #                 splits = splitter.split_documents(docs)
                
# # #                 # Assign default public access to uploaded files for consistency
# # #                 for split in splits:
# # #                     split.metadata["access"] = "public"

# # #                 vectorstore = Chroma.from_documents(
# # #                     documents=splits, 
# # #                     embedding=embeddings_model,
# # #                     collection_name=f"temp_{file_hash}"
# # #                 )
                
# # #                 if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
# # #                        temp_vector_stores_cache.pop(next(iter(temp_vector_stores_cache)))
# # #                 temp_vector_stores_cache[file_hash] = vectorstore
# # #             finally:
# # #                 if os.path.exists(tmp_path):
# # #                     os.remove(tmp_path)

# # #         # Temporary files are always considered 'user' accessible (they just uploaded it)
# # #         # Note: We pass role="admin" here so that the RAG filter doesn't block the file the user literally just uploaded.
# # #         data = await get_rag_response(query, vectorstore, str(request.base_url), role="admin")
# # #         return ChatResponse(query=query, **data)

# # #     except Exception as e:
# # #         print(f"Endpoint Error: {e}")
# # #         raise HTTPException(500, f"Processing failed: {str(e)}")

# # # if __name__ == "__main__":
# # #     uvicorn.run("app:app", host="127.0.0.1", port=8000, reload=True)

# # import os
# # import shutil
# # import asyncio
# # import hashlib
# # import tempfile
# # import uvicorn
# # import torch
# # import glob
# # import re
# # from concurrent.futures import ThreadPoolExecutor
# # from typing import Optional, List, Dict, Any

# # # --- API & FASTAPI IMPORTS ---
# # from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
# # from fastapi.middleware.cors import CORSMiddleware
# # from fastapi.responses import JSONResponse
# # from fastapi.staticfiles import StaticFiles
# # from pydantic import BaseModel

# # # --- LANGCHAIN & AI IMPORTS ---
# # from langchain_core.documents import Document
# # from langchain.text_splitter import RecursiveCharacterTextSplitter
# # from langchain_huggingface import HuggingFaceEmbeddings
# # from langchain_chroma import Chroma
# # from langchain.prompts import PromptTemplate
# # from langchain_core.output_parsers import StrOutputParser
# # from langchain_ollama import OllamaLLM
# # from langchain_core.vectorstores import VectorStore

# # # --- OCR & VISION IMPORTS ---
# # import pytesseract
# # from pdf2image import convert_from_path
# # from PIL import Image, ImageOps

# # # --- DOCX IMPORT ---
# # try:
# #     from docx import Document as DocxDocument
# #     DOCX_AVAILABLE = True
# # except ImportError:
# #     DOCX_AVAILABLE = False
# #     print("WARNING: 'python-docx' not found. .docx support disabled.")

# # # --- PPTX IMPORT ---
# # try:
# #     from pptx import Presentation
# #     PPTX_AVAILABLE = True
# # except ImportError:
# #     PPTX_AVAILABLE = False
# #     print("WARNING: 'python-pptx' not found. .pptx support disabled.")

# # # --- TTS IMPORT ---
# # try:
# #     from TTS.api import TTS
# #     TTS_AVAILABLE = True
# # except ImportError:
# #     print("WARNING: 'TTS' library not found. Audio generation will be disabled.")
# #     TTS_AVAILABLE = False


# # # ===================================================================
# # # 1. CONFIGURATION & PATH MANAGEMENT
# # # ===================================================================

# # class AgentConfig:
# #     # 1. Tesseract Path
# #     TESSERACT_CMD = shutil.which("tesseract") or r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
# #     # 2. Poppler Path (For PDF -> Image)
# #     POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
    
# #     # 3. Storage Paths
# #     CHROMA_DB_PATH = "./chroma_db_folder"
# #     KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\test") 
# #     AUDIO_DIR = "./static_audio"

# #     # 4. Model Settings
# #     LLM_MODEL = "phi3"
# #     EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
# #     # 5. Supported Extensions
# #     ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg'}

# # # Apply Configuration
# # os.environ['HF_HUB_OFFLINE'] = '1'
# # os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)
# # pytesseract.pytesseract.tesseract_cmd = AgentConfig.TESSERACT_CMD

# # # Verify Critical Dependencies
# # if not os.path.exists(AgentConfig.TESSERACT_CMD):
# #     print(f"CRITICAL WARNING: Tesseract not found at {AgentConfig.TESSERACT_CMD}. OCR will fail.")
# # if not os.path.exists(AgentConfig.POPPLER_PATH):
# #     print(f"CRITICAL WARNING: Poppler not found at {AgentConfig.POPPLER_PATH}. PDF processing will fail.")


# # # ===================================================================
# # # 2. GLOBAL RESOURCES & THREAD POOL
# # # ===================================================================
# # thread_pool = ThreadPoolExecutor(max_workers=4)
# # temp_vector_stores_cache = {}
# # MAX_CACHE_SIZE = 10 

# # # --- LOAD MODELS ---
# # print("\n--- Loading Models ---")

# # # 1. Embeddings
# # print(f"Loading Embedding Model ({AgentConfig.EMBEDDING_MODEL})...")
# # try:
# #     device = 'cuda' if torch.cuda.is_available() else 'cpu'
# #     embeddings_model = HuggingFaceEmbeddings(
# #         model_name=AgentConfig.EMBEDDING_MODEL,
# #         model_kwargs={'device': device},
# #         encode_kwargs={'normalize_embeddings': False}
# #     )
# #     print("✓ Embeddings Loaded")
# # except Exception as e:
# #     print(f"❌ Failed to load embeddings: {e}")
# #     exit(1)

# # # 2. LLM
# # print(f"Connecting to Ollama ({AgentConfig.LLM_MODEL})...")
# # try:
# #     llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
# #     print("✓ Ollama Connected")
# # except Exception as e:
# #     print(f"❌ Failed to connect to Ollama: {e}")

# # # 3. TTS
# # tts_model = None
# # if TTS_AVAILABLE:
# #     print("Loading TTS Model...")
# #     try:
# #         tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC")
# #         if torch.cuda.is_available():
# #             tts_model.to("cuda")
# #         print("✓ TTS Model Loaded")
# #     except Exception as e:
# #         print(f"⚠ TTS Load Error: {e}. Audio features disabled.")
# #         tts_model = None


# # # ===================================================================
# # # 3. IMAGE PRE-PROCESSING
# # # ===================================================================

# # def improve_image_quality(image: Image.Image) -> Image.Image:
# #     """
# #     Pre-processes an image to make it easier for Tesseract to read.
# #     """
# #     try:
# #         image = image.convert('L')
# #         width, height = image.size
# #         if width < 1000:
# #             scale_factor = 2
# #             image = image.resize((width * scale_factor, height * scale_factor), Image.Resampling.LANCZOS)
# #         image = image.point(lambda x: 0 if x < 140 else 255, '1')
# #         return image
# #     except Exception as e:
# #         print(f"Image preprocessing warning: {e}")
# #         return image 

# # def clean_ocr_text(text: str) -> str:
# #     """Removes garbage characters often produced by OCR."""
# #     text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
# #     text = re.sub(r'\s+([|.,;:])\s+', r'\1 ', text)
# #     return text.strip()


# # # ===================================================================
# # # 4. CORE PROCESSING FUNCTION
# # # ===================================================================

# # def _run_processing_sync(file_path: str, file_name: str) -> List[Document]:
# #     """Blocking function to process PDF, DOCX, PPTX, and Images."""
# #     extracted_docs = []
# #     file_lower = file_name.lower()
# #     ext = os.path.splitext(file_lower)[1]

# #     try:
# #         # --- PDF PROCESSING ---
# #         if ext == '.pdf':
# #             images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
# #             for i, image in enumerate(images):
# #                 processed_img = improve_image_quality(image)
# #                 text = pytesseract.image_to_string(processed_img, lang='eng')
# #                 clean = clean_ocr_text(text)
                
# #                 if clean:
# #                     extracted_docs.append(Document(
# #                         page_content=clean,
# #                         metadata={'source': file_name, 'page': i + 1}
# #                     ))

# #         # --- DOCX PROCESSING ---
# #         elif ext == '.docx':
# #             if not DOCX_AVAILABLE: raise ImportError("python-docx not installed.")
# #             doc = DocxDocument(file_path)
# #             full_text = []
# #             for para in doc.paragraphs:
# #                 if para.text.strip(): full_text.append(para.text)
# #             for table in doc.tables:
# #                 table_content = []
# #                 for row in table.rows:
# #                     row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
# #                     if row_cells: table_content.append(" | ".join(row_cells))
# #                 if table_content: full_text.append("\n[TABLE DATA]\n" + "\n".join(table_content) + "\n")

# #             combined_text = "\n".join(full_text)
# #             if combined_text.strip():
# #                 extracted_docs.append(Document(page_content=combined_text, metadata={'source': file_name, 'page': 1}))

# #         # --- PPTX PROCESSING ---
# #         elif ext == '.pptx':
# #             if not PPTX_AVAILABLE: raise ImportError("python-pptx not installed.")
# #             prs = Presentation(file_path)
# #             for i, slide in enumerate(prs.slides):
# #                 slide_text = []
# #                 for shape in slide.shapes:
# #                     if hasattr(shape, "text") and shape.text.strip():
# #                         slide_text.append(shape.text)
# #                     if shape.has_table:
# #                         table_rows = []
# #                         for row in shape.table.rows:
# #                             row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
# #                             if row_cells: table_rows.append(" | ".join(row_cells))
# #                         if table_rows: slide_text.append("\n[SLIDE TABLE]\n" + "\n".join(table_rows))

# #                 full_text = "\n".join(slide_text)
# #                 if full_text.strip():
# #                     extracted_docs.append(Document(page_content=full_text, metadata={'source': file_name, 'page': i + 1}))

# #         # --- IMAGE PROCESSING ---
# #         elif ext in ['.png', '.jpg', '.jpeg']:
# #             image = Image.open(file_path)
# #             processed_img = improve_image_quality(image)
# #             custom_config = r'--oem 3 --psm 3'
# #             text = pytesseract.image_to_string(processed_img, lang='eng', config=custom_config)
# #             clean = clean_ocr_text(text)
# #             if clean:
# #                 extracted_docs.append(Document(page_content=clean, metadata={'source': file_name, 'page': 1}))

# #     except Exception as e:
# #         print(f"Error processing {file_name}: {e}")
    
# #     return extracted_docs

# # async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
# #     loop = asyncio.get_running_loop()
# #     return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


# # def _run_tts_sync(text: str, file_path: str):
# #     if tts_model:
# #         tts_model.tts_to_file(text=text, file_path=file_path)

# # async def generate_audio_async(text: str, base_url: str) -> Optional[str]:
# #     if not tts_model or not text.strip():
# #         return None
# #     try:
# #         text_hash = hashlib.md5(text.encode()).hexdigest()
# #         filename = f"{text_hash}.wav"
# #         file_path = os.path.join(AgentConfig.AUDIO_DIR, filename)
# #         if not os.path.exists(file_path):
# #             loop = asyncio.get_running_loop()
# #             await loop.run_in_executor(thread_pool, _run_tts_sync, text, file_path)
# #         return f"{base_url}audio/{filename}"
# #     except Exception:
# #         return None


# # # ===================================================================
# # # 5. RAG PIPELINE (STRICT SECURITY & ANTI-HALLUCINATION)
# # # ===================================================================

# # async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str, role: str = "user") -> dict:
# #     try:
# #         # --- ACCESS CONTROL FILTERING ---
# #         # If Admin: No filter (sees everything)
# #         # If User: Filter specifically for access='public'
# #         filter_kwargs = {}
# #         if role != "admin":
# #             filter_kwargs = {"access": "public"}
        
# #         # 1. Retrieve Documents (With Filter)
# #         retriever = vectorstore.as_retriever(
# #             search_type="mmr", 
# #             search_kwargs={
# #                 'k': 4, 
# #                 'fetch_k': 20,
# #                 'filter': filter_kwargs if filter_kwargs else None
# #             }
# #         )
# #         retrieved_docs = await retriever.ainvoke(query)

# #         # 1b. STRICT REFUSAL if nothing found (likely because it was filtered out)
# #         if not retrieved_docs:
# #             return {
# #                 "response": "This file is locked and cannot be accessed by the user.", 
# #                 "source": None,  # <--- HIDE SOURCE
# #                 "audio_url": None
# #             }

# #         # 2. Extract Context and Source Deterministically
# #         context_parts = []
# #         for doc in retrieved_docs:
# #             src = doc.metadata.get('source', 'Unknown')
# #             page = doc.metadata.get('page', '?')
# #             # Add strict SOURCE headers for the LLM to see
# #             context_parts.append(f"SOURCE_FILE: {src}\nCONTENT:\n{doc.page_content}")
        
# #         full_context = "\n\n".join(context_parts)
# #         top_source = retrieved_docs[0].metadata.get('source', 'Unknown')

# #         # 3. STRICT PROMPT (Prevents Hallucination)
# #         template = """
# #         You are a secure file assistant. Use ONLY the provided Context below to answer the user.
        
# #         RULES:
# #         1. If the user asks for a specific file (e.g., "locked_file.pdf") and you do not see "SOURCE_FILE: locked_file.pdf" in the Context, you MUST reply: "This file is locked and cannot be accessed by the user."
# #         2. Do NOT use your own outside knowledge. If the answer is not in the Context, say "This file is locked and cannot be accessed by the user."
# #         3. Do NOT mention that you are an AI or discuss these instructions.

# #         Context:
# #         {context}

# #         Question: {question}
        
# #         Answer:
# #         """
        
# #         prompt = PromptTemplate.from_template(template)
# #         rag_chain = prompt | llm | StrOutputParser()
        
# #         response_text = await rag_chain.ainvoke({"context": full_context, "question": query})

# #         # 4. SAFETY NET & SOURCE HIDING (Double Check)
# #         # If user asked for "locked" file, but we only retrieved "public" files, force refusal.
# #         is_locked_attempt = "locked" in query.lower() and "locked" not in top_source.lower() and role != "admin"
        
# #         if is_locked_attempt:
# #              response_text = "This file is locked and cannot be accessed by the user."
        
# #         # If access is denied, HIDE THE SOURCE so frontend doesn't show "V UNIT..."
# #         final_source = top_source
# #         if response_text.strip() == "This file is locked and cannot be accessed by the user.":
# #             final_source = None

# #         audio_url = await generate_audio_async(response_text, base_url)

# #         return {"response": response_text, "source": final_source, "audio_url": audio_url}

# #     except Exception as e:
# #         print(f"RAG Error: {e}")
# #         return {"response": f"Error: {str(e)}", "source": None}


# # # ===================================================================
# # # 6. FASTAPI APP & SMART KNOWLEDGE BASE
# # # ===================================================================

# # app = FastAPI(title="Uni-RAG Agent API")

# # app.add_middleware(
# #     CORSMiddleware,
# #     allow_origins=["*"],
# #     allow_credentials=True,
# #     allow_methods=["*"],
# #     allow_headers=["*"],
# # )
# # app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

# # folder_vector_store = None

# # def init_knowledge_base():
# #     """Scans RAG-Base, tags files as admin/public based on name, and saves to ChromaDB."""
# #     global folder_vector_store
    
# #     # 1. Initialize persistent DB
# #     folder_vector_store = Chroma(
# #         persist_directory=AgentConfig.CHROMA_DB_PATH,
# #         embedding_function=embeddings_model
# #     )
    
# #     if not os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
# #         print(f"Warning: KB path {AgentConfig.KNOWLEDGE_BASE_PATH} does not exist.")
# #         return

# #     # 2. Get list of files already in DB
# #     print("--- Scanning Knowledge Base for New Files ---")
# #     existing_data = folder_vector_store.get()
# #     existing_sources = set()
# #     if existing_data and 'metadatas' in existing_data:
# #         for meta in existing_data['metadatas']:
# #             if meta and 'source' in meta:
# #                 existing_sources.add(meta['source'])
    
# #     print(f"Found {len(existing_sources)} existing files in DB.")

# #     # 3. Scan directory for files
# #     files_to_process = []
# #     try:
# #         all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
# #         for f in all_files:
# #             ext = os.path.splitext(f)[1].lower()
# #             if ext in AgentConfig.ALLOWED_EXTENSIONS:
# #                 if f not in existing_sources:
# #                     files_to_process.append(f)
# #     except Exception as e:
# #         print(f"Error reading RAG-Base directory: {e}")

# #     # 4. Process only NEW files
# #     if not files_to_process:
# #         print("✓ Knowledge Base is up to date.")
# #     else:
# #         print(f"Found {len(files_to_process)} new files to ingest.")
# #         new_docs = []
# #         splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

# #         for filename in files_to_process:
# #             file_path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, filename)
# #             print(f"Processing: {filename}...")
            
# #             # --- AUTOMATIC SECURITY TAGGING ---
# #             is_locked = filename.lower().startswith("locked_")
# #             access_level = "admin" if is_locked else "public"
# #             # ----------------------------------

# #             docs = _run_processing_sync(file_path, filename) 
            
# #             if docs:
# #                 # Inject access level into metadata for every chunk
# #                 for doc in docs:
# #                     doc.metadata["access"] = access_level

# #                 split_docs = splitter.split_documents(docs)
# #                 new_docs.extend(split_docs)
# #             else:
# #                 print(f"  Warning: No text extracted from {filename}")

# #         if new_docs:
# #             print(f"Adding {len(new_docs)} chunks to ChromaDB...")
# #             folder_vector_store.add_documents(new_docs)
# #             print("✓ Ingestion Complete!")
# #         else:
# #             print("No valid content found in new files.")


# # # Trigger Smart Ingestion on Startup
# # init_knowledge_base()


# # class ChatRequest(BaseModel):
# #     query: str
# #     role: str = "user"  # Default role is user

# # class ChatResponse(BaseModel):
# #     query: str
# #     response: str
# #     source: Optional[str] = None
# #     audio_url: Optional[str] = None


# # @app.post("/chat-folder", response_model=ChatResponse)
# # async def chat_with_folder(request: ChatRequest, req: Request):
# #     if not folder_vector_store:
# #         raise HTTPException(500, "Knowledge Base not loaded.")
    
# #     # Pass the role to the RAG function
# #     data = await get_rag_response(
# #         query=request.query, 
# #         vectorstore=folder_vector_store, 
# #         base_url=str(req.base_url),
# #         role=request.role
# #     )
    
# #     return ChatResponse(query=request.query, **data)


# # @app.post("/chat-file", response_model=ChatResponse)
# # async def chat_with_single_file(
# #     request: Request,
# #     query: str = Form(...),
# #     file: UploadFile = File(...)
# # ):
# #     try:
# #         suffix = os.path.splitext(file.filename)[1].lower()
# #         if suffix not in AgentConfig.ALLOWED_EXTENSIONS:
# #             raise HTTPException(400, f"Invalid file. Allowed: {AgentConfig.ALLOWED_EXTENSIONS}")

# #         content = await file.read()
# #         file_hash = hashlib.md5(content).hexdigest()
        
# #         if file_hash in temp_vector_stores_cache:
# #             vectorstore = temp_vector_stores_cache[file_hash]
# #         else:
# #             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
# #                 tmp.write(content)
# #                 tmp_path = tmp.name
            
# #             try:
# #                 docs = await extract_text_async(tmp_path, file.filename)
                
# #                 if not docs:
# #                      raise HTTPException(400, "Could not extract text from file.")

# #                 splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
# #                 splits = splitter.split_documents(docs)
                
# #                 # Assign default public access to uploaded files for consistency
# #                 for split in splits:
# #                     split.metadata["access"] = "public"

# #                 vectorstore = Chroma.from_documents(
# #                     documents=splits, 
# #                     embedding=embeddings_model,
# #                     collection_name=f"temp_{file_hash}"
# #                 )
                
# #                 if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
# #                        temp_vector_stores_cache.pop(next(iter(temp_vector_stores_cache)))
# #                 temp_vector_stores_cache[file_hash] = vectorstore
# #             finally:
# #                 if os.path.exists(tmp_path):
# #                     os.remove(tmp_path)

# #         # Temporary files are always considered 'user' accessible (they just uploaded it)
# #         # Note: We pass role="admin" here so that the RAG filter doesn't block the file the user literally just uploaded.
# #         data = await get_rag_response(query, vectorstore, str(request.base_url), role="admin")
# #         return ChatResponse(query=query, **data)

# #     except Exception as e:
# #         print(f"Endpoint Error: {e}")
# #         raise HTTPException(500, f"Processing failed: {str(e)}")

# # if __name__ == "__main__":
# #     uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)









# import os
# import shutil
# import asyncio
# import hashlib
# import tempfile
# import uvicorn
# import torch
# import glob
# import re
# from concurrent.futures import ThreadPoolExecutor
# from typing import Optional, List, Dict, Any

# # --- API & FASTAPI IMPORTS ---
# from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
# from fastapi.middleware.cors import CORSMiddleware
# from fastapi.responses import JSONResponse
# from fastapi.staticfiles import StaticFiles
# from pydantic import BaseModel

# # --- LANGCHAIN & AI IMPORTS ---
# from langchain_core.documents import Document
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_chroma import Chroma
# from langchain.prompts import PromptTemplate
# from langchain_core.output_parsers import StrOutputParser
# from langchain_ollama import OllamaLLM
# from langchain_core.vectorstores import VectorStore

# # --- OCR & VISION IMPORTS ---
# import pytesseract
# from pdf2image import convert_from_path
# from PIL import Image, ImageOps

# # --- DOCX IMPORT ---
# try:
#     from docx import Document as DocxDocument
#     DOCX_AVAILABLE = True
# except ImportError:
#     DOCX_AVAILABLE = False
#     print("WARNING: 'python-docx' not found. .docx support disabled.")

# # --- PPTX IMPORT ---
# try:
#     from pptx import Presentation
#     PPTX_AVAILABLE = True
# except ImportError:
#     PPTX_AVAILABLE = False
#     print("WARNING: 'python-pptx' not found. .pptx support disabled.")

# # --- TTS IMPORT ---
# try:
#     from TTS.api import TTS
#     TTS_AVAILABLE = True
# except ImportError:
#     print("WARNING: 'TTS' library not found. Audio generation will be disabled.")
#     TTS_AVAILABLE = False


# # ===================================================================
# # 1. CONFIGURATION & PATH MANAGEMENT
# # ===================================================================

# class AgentConfig:
#     # 1. Tesseract Path
#     TESSERACT_CMD = shutil.which("tesseract") or r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
#     # 2. Poppler Path (For PDF -> Image)
#     POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
    
#     # 3. Storage Paths
#     CHROMA_DB_PATH = "./chroma_db_folder"
#     KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\test") 
#     AUDIO_DIR = "./static_audio"

#     # 4. Model Settings
#     LLM_MODEL = "phi3"
#     EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
#     # 5. Supported Extensions
#     ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg'}

# # Apply Configuration
# os.environ['HF_HUB_OFFLINE'] = '1'
# os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)
# pytesseract.pytesseract.tesseract_cmd = AgentConfig.TESSERACT_CMD

# # Verify Critical Dependencies
# if not os.path.exists(AgentConfig.TESSERACT_CMD):
#     print(f"CRITICAL WARNING: Tesseract not found at {AgentConfig.TESSERACT_CMD}. OCR will fail.")
# if not os.path.exists(AgentConfig.POPPLER_PATH):
#     print(f"CRITICAL WARNING: Poppler not found at {AgentConfig.POPPLER_PATH}. PDF processing will fail.")


# # ===================================================================
# # 2. GLOBAL RESOURCES & THREAD POOL
# # ===================================================================
# thread_pool = ThreadPoolExecutor(max_workers=4)
# temp_vector_stores_cache = {}
# MAX_CACHE_SIZE = 10 

# # --- LOAD MODELS ---
# print("\n--- Loading Models ---")

# # 1. Embeddings
# print(f"Loading Embedding Model ({AgentConfig.EMBEDDING_MODEL})...")
# try:
#     device = 'cuda' if torch.cuda.is_available() else 'cpu'
#     embeddings_model = HuggingFaceEmbeddings(
#         model_name=AgentConfig.EMBEDDING_MODEL,
#         model_kwargs={'device': device},
#         encode_kwargs={'normalize_embeddings': False}
#     )
#     print("✓ Embeddings Loaded")
# except Exception as e:
#     print(f"❌ Failed to load embeddings: {e}")
#     exit(1)

# # 2. LLM
# print(f"Connecting to Ollama ({AgentConfig.LLM_MODEL})...")
# try:
#     llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
#     print("✓ Ollama Connected")
# except Exception as e:
#     print(f"❌ Failed to connect to Ollama: {e}")

# # 3. TTS
# tts_model = None
# if TTS_AVAILABLE:
#     print("Loading TTS Model...")
#     try:
#         tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC")
#         if torch.cuda.is_available():
#             tts_model.to("cuda")
#         print("✓ TTS Model Loaded")
#     except Exception as e:
#         print(f"⚠ TTS Load Error: {e}. Audio features disabled.")
#         tts_model = None


# # ===================================================================
# # 3. IMAGE PRE-PROCESSING
# # ===================================================================

# def improve_image_quality(image: Image.Image) -> Image.Image:
#     """
#     Pre-processes an image to make it easier for Tesseract to read.
#     """
#     try:
#         image = image.convert('L')
#         width, height = image.size
#         if width < 1000:
#             scale_factor = 2
#             image = image.resize((width * scale_factor, height * scale_factor), Image.Resampling.LANCZOS)
#         image = image.point(lambda x: 0 if x < 140 else 255, '1')
#         return image
#     except Exception as e:
#         print(f"Image preprocessing warning: {e}")
#         return image 

# def clean_ocr_text(text: str) -> str:
#     """Removes garbage characters often produced by OCR."""
#     text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
#     text = re.sub(r'\s+([|.,;:])\s+', r'\1 ', text)
#     return text.strip()


# # ===================================================================
# # 4. CORE PROCESSING FUNCTION
# # ===================================================================

# def _run_processing_sync(file_path: str, clean_name: str) -> List[Document]:
#     """Blocking function to process PDF, DOCX, PPTX, and Images."""
#     extracted_docs = []
#     file_lower = file_path.lower() # Use path for extension check, but clean_name for metadata
#     ext = os.path.splitext(file_lower)[1]

#     try:
#         # --- PDF PROCESSING ---
#         if ext == '.pdf':
#             images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
#             for i, image in enumerate(images):
#                 processed_img = improve_image_quality(image)
#                 text = pytesseract.image_to_string(processed_img, lang='eng')
#                 clean = clean_ocr_text(text)
                
#                 if clean:
#                     extracted_docs.append(Document(
#                         page_content=clean,
#                         metadata={'source': clean_name, 'page': i + 1}
#                     ))

#         # --- DOCX PROCESSING ---
#         elif ext == '.docx':
#             if not DOCX_AVAILABLE: raise ImportError("python-docx not installed.")
#             doc = DocxDocument(file_path)
#             full_text = []
#             for para in doc.paragraphs:
#                 if para.text.strip(): full_text.append(para.text)
#             for table in doc.tables:
#                 table_content = []
#                 for row in table.rows:
#                     row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
#                     if row_cells: table_content.append(" | ".join(row_cells))
#                 if table_content: full_text.append("\n[TABLE DATA]\n" + "\n".join(table_content) + "\n")

#             combined_text = "\n".join(full_text)
#             if combined_text.strip():
#                 extracted_docs.append(Document(page_content=combined_text, metadata={'source': clean_name, 'page': 1}))

#         # --- PPTX PROCESSING ---
#         elif ext == '.pptx':
#             if not PPTX_AVAILABLE: raise ImportError("python-pptx not installed.")
#             prs = Presentation(file_path)
#             for i, slide in enumerate(prs.slides):
#                 slide_text = []
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text") and shape.text.strip():
#                         slide_text.append(shape.text)
#                     if shape.has_table:
#                         table_rows = []
#                         for row in shape.table.rows:
#                             row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
#                             if row_cells: table_rows.append(" | ".join(row_cells))
#                         if table_rows: slide_text.append("\n[SLIDE TABLE]\n" + "\n".join(table_rows))

#                 full_text = "\n".join(slide_text)
#                 if full_text.strip():
#                     extracted_docs.append(Document(page_content=full_text, metadata={'source': clean_name, 'page': i + 1}))

#         # --- IMAGE PROCESSING ---
#         elif ext in ['.png', '.jpg', '.jpeg']:
#             image = Image.open(file_path)
#             processed_img = improve_image_quality(image)
#             custom_config = r'--oem 3 --psm 3'
#             text = pytesseract.image_to_string(processed_img, lang='eng', config=custom_config)
#             clean = clean_ocr_text(text)
#             if clean:
#                 extracted_docs.append(Document(page_content=clean, metadata={'source': clean_name, 'page': 1}))

#     except Exception as e:
#         print(f"Error processing {clean_name}: {e}")
    
#     return extracted_docs

# async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
#     loop = asyncio.get_running_loop()
#     return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


# def _run_tts_sync(text: str, file_path: str):
#     if tts_model:
#         tts_model.tts_to_file(text=text, file_path=file_path)

# async def generate_audio_async(text: str, base_url: str) -> Optional[str]:
#     if not tts_model or not text.strip():
#         return None
#     try:
#         text_hash = hashlib.md5(text.encode()).hexdigest()
#         filename = f"{text_hash}.wav"
#         file_path = os.path.join(AgentConfig.AUDIO_DIR, filename)
#         if not os.path.exists(file_path):
#             loop = asyncio.get_running_loop()
#             await loop.run_in_executor(thread_pool, _run_tts_sync, text, file_path)
#         return f"{base_url}audio/{filename}"
#     except Exception:
#         return None


# # ===================================================================
# # 5. RAG PIPELINE (STRICT SECURITY & ANTI-HALLUCINATION)
# # ===================================================================

# async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str, role: str = "user") -> dict:
#     try:
#         # --- ACCESS CONTROL FILTERING ---
#         filter_kwargs = {}
#         if role != "admin":
#             filter_kwargs = {"access": "public"}
        
#         # 1. Retrieve Documents
#         retriever = vectorstore.as_retriever(
#             search_type="mmr", 
#             search_kwargs={
#                 'k': 4, 
#                 'fetch_k': 20,
#                 'filter': filter_kwargs if filter_kwargs else None
#             }
#         )
#         retrieved_docs = await retriever.ainvoke(query)

#         # 1b. STRICT REFUSAL if nothing found
#         if not retrieved_docs:
#             return {
#                 "response": "This file is locked and cannot be accessed by the user.", 
#                 "source": None, 
#                 "audio_url": None
#             }

#         # 2. Extract Context and Source
#         context_parts = []
#         for doc in retrieved_docs:
#             src = doc.metadata.get('source', 'Unknown')
#             context_parts.append(f"SOURCE_FILE: {src}\nCONTENT:\n{doc.page_content}")
        
#         full_context = "\n\n".join(context_parts)
#         top_source = retrieved_docs[0].metadata.get('source', 'Unknown')

#         # 3. STRICT PROMPT
#         # Note: We removed the explicit "locked_file" check from the prompt instructions 
#         # because the LLM will now see clean names (e.g. "salary.pdf") in the context.
#         # The database filter handles the security. If the file is locked, the context is empty, and it returns refusal.
#         template = """
#         You are a secure file assistant. Use ONLY the provided Context below to answer the user.
        
#         RULES:
#         1. If the answer is not in the Context, you MUST reply: "This file is locked and cannot be accessed by the user."
#         2. Do NOT use your own outside knowledge. 
#         3. Do NOT mention that you are an AI.

#         Context:
#         {context}

#         Question: {question}
        
#         Answer:
#         """
        
#         prompt = PromptTemplate.from_template(template)
#         rag_chain = prompt | llm | StrOutputParser()
        
#         response_text = await rag_chain.ainvoke({"context": full_context, "question": query})

#         # 4. SAFETY NET & SOURCE HIDING
#         # If access is denied, HIDE THE SOURCE so frontend doesn't show "V UNIT..."
#         final_source = top_source
#         if response_text.strip() == "This file is locked and cannot be accessed by the user.":
#             final_source = None

#         audio_url = await generate_audio_async(response_text, base_url)

#         return {"response": response_text, "source": final_source, "audio_url": audio_url}

#     except Exception as e:
#         print(f"RAG Error: {e}")
#         return {"response": f"Error: {str(e)}", "source": None}


# # ===================================================================
# # 6. FASTAPI APP & SMART KNOWLEDGE BASE
# # ===================================================================

# app = FastAPI(title="Uni-RAG Agent API")

# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=["*"],
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"],
# )
# app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

# folder_vector_store = None

# def init_knowledge_base():
#     """Scans RAG-Base, strips 'locked_' prefix for storage, sets metadata, and saves to ChromaDB."""
#     global folder_vector_store
    
#     # 1. Initialize persistent DB
#     folder_vector_store = Chroma(
#         persist_directory=AgentConfig.CHROMA_DB_PATH,
#         embedding_function=embeddings_model
#     )
    
#     if not os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
#         print(f"Warning: KB path {AgentConfig.KNOWLEDGE_BASE_PATH} does not exist.")
#         return

#     # 2. Get list of files already in DB (These will be the CLEAN names now)
#     print("--- Scanning Knowledge Base for New Files ---")
#     existing_data = folder_vector_store.get()
#     existing_sources = set()
#     if existing_data and 'metadatas' in existing_data:
#         for meta in existing_data['metadatas']:
#             if meta and 'source' in meta:
#                 existing_sources.add(meta['source'])
    
#     print(f"Found {len(existing_sources)} existing files in DB.")

#     # 3. Scan directory for files
#     files_to_process = []
#     try:
#         all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
#         for f in all_files:
#             ext = os.path.splitext(f)[1].lower()
#             if ext in AgentConfig.ALLOWED_EXTENSIONS:
#                 # --- CALCULATE CLEAN NAME TO CHECK EXISTENCE ---
#                 # We must check if the *clean* name exists in DB to prevent re-ingesting locked files repeatedly
#                 is_locked = f.lower().startswith("locked_")
#                 clean_name = f[7:] if is_locked else f
                
#                 if clean_name not in existing_sources:
#                     files_to_process.append(f)
                    
#     except Exception as e:
#         print(f"Error reading RAG-Base directory: {e}")

#     # 4. Process only NEW files
#     if not files_to_process:
#         print("✓ Knowledge Base is up to date.")
#     else:
#         print(f"Found {len(files_to_process)} new files to ingest.")
#         new_docs = []
#         splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

#         for filename in files_to_process:
#             file_path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, filename)
            
#             # --- AUTOMATIC SECURITY & NAME CLEANING ---
#             is_locked = filename.lower().startswith("locked_")
#             access_level = "admin" if is_locked else "public"
            
#             # Determine the name we want to store in ChromaDB (stripping prefix)
#             clean_filename = filename[7:] if is_locked else filename
#             # ------------------------------------------

#             print(f"Processing: {filename} -> Storing as: {clean_filename} [{access_level}]")

#             # Pass the CLEAN name to the processing function
#             docs = _run_processing_sync(file_path, clean_filename) 
            
#             if docs:
#                 # Inject access level into metadata for every chunk
#                 for doc in docs:
#                     doc.metadata["access"] = access_level
#                     # Ensure source is the clean name (redundant but safe)
#                     doc.metadata["source"] = clean_filename

#                 split_docs = splitter.split_documents(docs)
#                 new_docs.extend(split_docs)
#             else:
#                 print(f"  Warning: No text extracted from {filename}")

#         if new_docs:
#             print(f"Adding {len(new_docs)} chunks to ChromaDB...")
#             folder_vector_store.add_documents(new_docs)
#             print("✓ Ingestion Complete!")
#         else:
#             print("No valid content found in new files.")


# # Trigger Smart Ingestion on Startup
# init_knowledge_base()


# class ChatRequest(BaseModel):
#     query: str
#     role: str = "admin"

# class ChatResponse(BaseModel):
#     query: str
#     response: str
#     source: Optional[str] = None
#     audio_url: Optional[str] = None


# @app.post("/chat-folder", response_model=ChatResponse)
# async def chat_with_folder(request: ChatRequest, req: Request):
#     if not folder_vector_store:
#         raise HTTPException(500, "Knowledge Base not loaded.")
    
#     data = await get_rag_response(
#         query=request.query, 
#         vectorstore=folder_vector_store, 
#         base_url=str(req.base_url),
#         role=request.role
#     )
    
#     return ChatResponse(query=request.query, **data)


# @app.post("/chat-file", response_model=ChatResponse)
# async def chat_with_single_file(
#     request: Request,
#     query: str = Form(...),
#     file: UploadFile = File(...)
# ):
#     try:
#         suffix = os.path.splitext(file.filename)[1].lower()
#         if suffix not in AgentConfig.ALLOWED_EXTENSIONS:
#             raise HTTPException(400, f"Invalid file. Allowed: {AgentConfig.ALLOWED_EXTENSIONS}")

#         content = await file.read()
#         file_hash = hashlib.md5(content).hexdigest()
        
#         if file_hash in temp_vector_stores_cache:
#             vectorstore = temp_vector_stores_cache[file_hash]
#         else:
#             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
#                 tmp.write(content)
#                 tmp_path = tmp.name
            
#             try:
#                 # For uploaded files, we usually just use the filename as is
#                 docs = await extract_text_async(tmp_path, file.filename)
                
#                 if not docs:
#                      raise HTTPException(400, "Could not extract text from file.")

#                 splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#                 splits = splitter.split_documents(docs)
                
#                 for split in splits:
#                     split.metadata["access"] = "public"

#                 vectorstore = Chroma.from_documents(
#                     documents=splits, 
#                     embedding=embeddings_model,
#                     collection_name=f"temp_{file_hash}"
#                 )
                
#                 if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
#                        temp_vector_stores_cache.pop(next(iter(temp_vector_stores_cache)))
#                 temp_vector_stores_cache[file_hash] = vectorstore
#             finally:
#                 if os.path.exists(tmp_path):
#                     os.remove(tmp_path)

#         data = await get_rag_response(query, vectorstore, str(request.base_url), role="admin")
#         return ChatResponse(query=query, **data)

#     except Exception as e:
#         print(f"Endpoint Error: {e}")
#         raise HTTPException(500, f"Processing failed: {str(e)}")

# if __name__ == "__main__":
#     uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)









































import os
import shutil
import asyncio
import hashlib
import tempfile
import uvicorn
import torch
import glob
import re
from concurrent.futures import ThreadPoolExecutor
from typing import Optional, List, Dict, Any

# --- API & FASTAPI IMPORTS ---
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

# --- LANGCHAIN & AI IMPORTS ---
from langchain_core.documents import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_ollama import OllamaLLM
from langchain_core.vectorstores import VectorStore

# --- OCR & VISION IMPORTS ---
import pytesseract
from pdf2image import convert_from_path
from PIL import Image, ImageOps

# --- DOCX IMPORT ---
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("WARNING: 'python-docx' not found. .docx support disabled.")

# --- PPTX IMPORT ---
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: 'python-pptx' not found. .pptx support disabled.")

# --- TTS IMPORT ---
try:
    from TTS.api import TTS
    TTS_AVAILABLE = True
except ImportError:
    print("WARNING: 'TTS' library not found. Audio generation will be disabled.")
    TTS_AVAILABLE = False


# ===================================================================
# 1. CONFIGURATION & PATH MANAGEMENT
# ===================================================================

class AgentConfig:
    # 1. Tesseract Path
    TESSERACT_CMD = shutil.which("tesseract") or r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
    # 2. Poppler Path (For PDF -> Image)
    POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
    
    # 3. Storage Paths
    CHROMA_DB_PATH = "./chroma_db_folder"
    KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\test") 
    AUDIO_DIR = "./static_audio"

    # 4. Model Settings
    LLM_MODEL = "phi3"
    EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
    # 5. Supported Extensions
    ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg'}

# Apply Configuration
os.environ['HF_HUB_OFFLINE'] = '1'
os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)
pytesseract.pytesseract.tesseract_cmd = AgentConfig.TESSERACT_CMD

# Verify Critical Dependencies
if not os.path.exists(AgentConfig.TESSERACT_CMD):
    print(f"CRITICAL WARNING: Tesseract not found at {AgentConfig.TESSERACT_CMD}. OCR will fail.")
if not os.path.exists(AgentConfig.POPPLER_PATH):
    print(f"CRITICAL WARNING: Poppler not found at {AgentConfig.POPPLER_PATH}. PDF processing will fail.")


# ===================================================================
# 2. GLOBAL RESOURCES & THREAD POOL
# ===================================================================
thread_pool = ThreadPoolExecutor(max_workers=4)
temp_vector_stores_cache = {}
MAX_CACHE_SIZE = 10 

# --- LOAD MODELS ---
print("\n--- Loading Models ---")

# 1. Embeddings
print(f"Loading Embedding Model ({AgentConfig.EMBEDDING_MODEL})...")
try:
    device = 'cuda' if torch.cuda.is_available() else 'cpu'
    embeddings_model = HuggingFaceEmbeddings(
        model_name=AgentConfig.EMBEDDING_MODEL,
        model_kwargs={'device': device},
        encode_kwargs={'normalize_embeddings': False}
    )
    print("✓ Embeddings Loaded")
except Exception as e:
    print(f"❌ Failed to load embeddings: {e}")
    exit(1)

# 2. LLM
print(f"Connecting to Ollama ({AgentConfig.LLM_MODEL})...")
try:
    llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
    print("✓ Ollama Connected")
except Exception as e:
    print(f"❌ Failed to connect to Ollama: {e}")

# 3. TTS
tts_model = None
if TTS_AVAILABLE:
    print("Loading TTS Model...")
    try:
        tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC")
        if torch.cuda.is_available():
            tts_model.to("cuda")
        print("✓ TTS Model Loaded")
    except Exception as e:
        print(f"⚠ TTS Load Error: {e}. Audio features disabled.")
        tts_model = None


# ===================================================================
# 3. IMAGE PRE-PROCESSING
# ===================================================================

def improve_image_quality(image: Image.Image) -> Image.Image:
    """Pre-processes an image to make it easier for Tesseract to read."""
    try:
        image = image.convert('L')
        width, height = image.size
        if width < 1000:
            scale_factor = 2
            image = image.resize((width * scale_factor, height * scale_factor), Image.Resampling.LANCZOS)
        image = image.point(lambda x: 0 if x < 140 else 255, '1')
        return image
    except Exception as e:
        print(f"Image preprocessing warning: {e}")
        return image 

def clean_ocr_text(text: str) -> str:
    """Removes garbage characters often produced by OCR."""
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    text = re.sub(r'\s+([|.,;:])\s+', r'\1 ', text)
    return text.strip()


# ===================================================================
# 4. CORE PROCESSING FUNCTION
# ===================================================================

def _run_processing_sync(file_path: str, file_name: str) -> List[Document]:
    """Blocking function to process PDF, DOCX, PPTX, and Images."""
    extracted_docs = []
    file_lower = file_name.lower()
    ext = os.path.splitext(file_lower)[1]

    try:
        # --- PDF PROCESSING ---
        if ext == '.pdf':
            images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
            for i, image in enumerate(images):
                processed_img = improve_image_quality(image)
                text = pytesseract.image_to_string(processed_img, lang='eng')
                clean = clean_ocr_text(text)
                if clean:
                    extracted_docs.append(Document(
                        page_content=clean,
                        metadata={'source': file_name, 'page': i + 1}
                    ))

        # --- DOCX PROCESSING ---
        elif ext == '.docx':
            if not DOCX_AVAILABLE: raise ImportError("python-docx not installed.")
            doc = DocxDocument(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip(): full_text.append(para.text)
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells: table_content.append(" | ".join(row_cells))
                if table_content: full_text.append("\n[TABLE DATA]\n" + "\n".join(table_content) + "\n")

            combined_text = "\n".join(full_text)
            if combined_text.strip():
                extracted_docs.append(Document(page_content=combined_text, metadata={'source': file_name, 'page': 1}))

        # --- PPTX PROCESSING ---
        elif ext == '.pptx':
            if not PPTX_AVAILABLE: raise ImportError("python-pptx not installed.")
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    if shape.has_table:
                        table_rows = []
                        for row in shape.table.rows:
                            row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
                            if row_cells: table_rows.append(" | ".join(row_cells))
                        if table_rows: slide_text.append("\n[SLIDE TABLE]\n" + "\n".join(table_rows))

                full_text = "\n".join(slide_text)
                if full_text.strip():
                    extracted_docs.append(Document(page_content=full_text, metadata={'source': file_name, 'page': i + 1}))

        # --- IMAGE PROCESSING ---
        elif ext in ['.png', '.jpg', '.jpeg']:
            image = Image.open(file_path)
            processed_img = improve_image_quality(image)
            custom_config = r'--oem 3 --psm 3'
            text = pytesseract.image_to_string(processed_img, lang='eng', config=custom_config)
            clean = clean_ocr_text(text)
            if clean:
                extracted_docs.append(Document(page_content=clean, metadata={'source': file_name, 'page': 1}))

    except Exception as e:
        print(f"Error processing {file_name}: {e}")
    
    return extracted_docs

async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


def _run_tts_sync(text: str, file_path: str):
    if tts_model:
        tts_model.tts_to_file(text=text, file_path=file_path)

async def generate_audio_async(text: str, base_url: str) -> Optional[str]:
    if not tts_model or not text.strip():
        return None
    try:
        text_hash = hashlib.md5(text.encode()).hexdigest()
        filename = f"{text_hash}.wav"
        file_path = os.path.join(AgentConfig.AUDIO_DIR, filename)
        if not os.path.exists(file_path):
            loop = asyncio.get_running_loop()
            await loop.run_in_executor(thread_pool, _run_tts_sync, text, file_path)
        return f"{base_url}audio/{filename}"
    except Exception:
        return None


# ===================================================================
# 5. RAG PIPELINE (STRICT SECURITY & ANTI-HALLUCINATION)
# ===================================================================

async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str, role: str = "user") -> dict:
    try:
        # --- ACCESS CONTROL FILTERING ---
        # 1. Admin: Can search everything.
        # 2. User: Can ONLY search documents tagged 'public'.
        filter_kwargs = {}
        if role != "admin":
            filter_kwargs = {"access": "public"}
        
        # 1. Retrieve Documents (Applying the Filter)
        retriever = vectorstore.as_retriever(
            search_type="mmr", 
            search_kwargs={
                'k': 4, 
                'fetch_k': 20,
                'filter': filter_kwargs if filter_kwargs else None
            }
        )
        retrieved_docs = await retriever.ainvoke(query)

        # 1b. STRICT REFUSAL: If filter blocked everything, stop here.
        if not retrieved_docs:
            return {
                "response": "This file is locked and cannot be accessed by the user.", 
                "source": None, 
                "audio_url": None
            }

        # 2. Extract Context and Source
        context_parts = []
        for doc in retrieved_docs:
            src = doc.metadata.get('source', 'Unknown')
            page = doc.metadata.get('page', '?')
            # Add strict SOURCE headers for the LLM to see
            context_parts.append(f"SOURCE_FILE: {src}\nCONTENT:\n{doc.page_content}")
        
        full_context = "\n\n".join(context_parts)
        top_source = retrieved_docs[0].metadata.get('source', 'Unknown')

        # 3. STRICT PROMPT (Anti-Hallucination)
        # We instruct the LLM: if the context is empty or doesn't match the user's request, DENY ACCESS.
        template = """
        You are a secure file assistant. Use ONLY the provided Context below to answer the user.
        
        RULES:
        1. If the answer is not explicitly in the Context, you MUST reply: "This file is locked and cannot be accessed by the user."
        2. Do NOT use your own outside knowledge.
        3. Do NOT mention that you are an AI.

        Context:
        {context}

        Question: {question}
        
        Answer:
        """
        
        prompt = PromptTemplate.from_template(template)
        rag_chain = prompt | llm | StrOutputParser()
        
        response_text = await rag_chain.ainvoke({"context": full_context, "question": query})

        # 4. SAFETY NET & SOURCE HIDING
        # If the LLM refuses access, we hide the source filename to prevent leaking existence of files.
        final_source = top_source
        if "locked and cannot be accessed" in response_text or response_text.strip() == "This file is locked and cannot be accessed by the user.":
            response_text = "This file is locked and cannot be accessed by the user."
            final_source = None

        audio_url = await generate_audio_async(response_text, base_url)

        return {"response": response_text, "source": final_source, "audio_url": audio_url}

    except Exception as e:
        print(f"RAG Error: {e}")
        return {"response": f"Error: {str(e)}", "source": None}


# ===================================================================
# 6. FASTAPI APP & SMART KNOWLEDGE BASE
# ===================================================================

app = FastAPI(title="Uni-RAG Agent API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

folder_vector_store = None

def init_knowledge_base():
    """Scans RAG-Base, tags files as admin/public based on name, and saves to ChromaDB."""
    global folder_vector_store
    
    # 1. Initialize persistent DB
    folder_vector_store = Chroma(
        persist_directory=AgentConfig.CHROMA_DB_PATH,
        embedding_function=embeddings_model
    )
    
    if not os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
        print(f"Warning: KB path {AgentConfig.KNOWLEDGE_BASE_PATH} does not exist.")
        return

    # 2. Get list of files already in DB
    print("--- Scanning Knowledge Base for New Files ---")
    existing_data = folder_vector_store.get()
    existing_sources = set()
    if existing_data and 'metadatas' in existing_data:
        for meta in existing_data['metadatas']:
            if meta and 'source' in meta:
                existing_sources.add(meta['source'])
    
    print(f"Found {len(existing_sources)} existing files in DB.")

    # 3. Scan directory for files
    files_to_process = []
    try:
        all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
        for f in all_files:
            ext = os.path.splitext(f)[1].lower()
            if ext in AgentConfig.ALLOWED_EXTENSIONS:
                if f not in existing_sources:
                    files_to_process.append(f)
    except Exception as e:
        print(f"Error reading RAG-Base directory: {e}")

    # 4. Process only NEW files
    if not files_to_process:
        print("✓ Knowledge Base is up to date.")
    else:
        print(f"Found {len(files_to_process)} new files to ingest.")
        new_docs = []
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

        for filename in files_to_process:
            file_path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, filename)
            print(f"Processing: {filename}...")
            
            # --- AUTOMATIC SECURITY TAGGING ---
            # If filename starts with "locked_", tag as ADMIN. Else PUBLIC.
            is_locked = filename.lower().startswith("locked_")
            access_level = "admin" if is_locked else "public"
            # ----------------------------------

            docs = _run_processing_sync(file_path, filename) 
            
            if docs:
                # Inject access level into metadata for every chunk
                for doc in docs:
                    doc.metadata["access"] = access_level

                split_docs = splitter.split_documents(docs)
                new_docs.extend(split_docs)
            else:
                print(f"  Warning: No text extracted from {filename}")

        if new_docs:
            print(f"Adding {len(new_docs)} chunks to ChromaDB...")
            folder_vector_store.add_documents(new_docs)
            print("✓ Ingestion Complete!")
        else:
            print("No valid content found in new files.")


# Trigger Smart Ingestion on Startup
init_knowledge_base()


class ChatRequest(BaseModel):
    query: str
    role: str = "admin"  # Default role is user. Frontend must send "admin" to see locked files.

class ChatResponse(BaseModel):
    query: str
    response: str
    source: Optional[str] = None
    audio_url: Optional[str] = None


# @app.post("/chat-folder", response_model=ChatResponse)
# async def chat_with_folder(request: ChatRequest, req: Request):
#     if not folder_vector_store:
#         raise HTTPException(500, "Knowledge Base not loaded.")
    
#     # Pass the role to the RAG function for filtering
#     data = await get_rag_response(
#         query=request.query, 
#         vectorstore=folder_vector_store, 
#         base_url=str(req.base_url),
#         role=request.role
#     )
    
#     return ChatResponse(query=request.query, **data)

@app.post("/chat-folder", response_model=ChatResponse)
async def chat_with_folder(request: ChatRequest, req: Request):
    if not folder_vector_store:
        raise HTTPException(500, "Knowledge Base not loaded.")
    
    # --- FORCE ADMIN MODE FOR TESTING ---
    # This ignores whatever the frontend sends and forces Admin access
    forced_role = "user"
    
    data = await get_rag_response(
        query=request.query, 
        vectorstore=folder_vector_store, 
        base_url=str(req.base_url),
        role=forced_role  # <--- Use the forced variable here
    )
    
    return ChatResponse(query=request.query, **data)


@app.post("/chat-file", response_model=ChatResponse)
async def chat_with_single_file(
    request: Request,
    query: str = Form(...),
    file: UploadFile = File(...)
):
    try:
        suffix = os.path.splitext(file.filename)[1].lower()
        if suffix not in AgentConfig.ALLOWED_EXTENSIONS:
            raise HTTPException(400, f"Invalid file. Allowed: {AgentConfig.ALLOWED_EXTENSIONS}")

        content = await file.read()
        file_hash = hashlib.md5(content).hexdigest()
        
        if file_hash in temp_vector_stores_cache:
            vectorstore = temp_vector_stores_cache[file_hash]
        else:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            
            try:
                docs = await extract_text_async(tmp_path, file.filename)
                
                if not docs:
                     raise HTTPException(400, "Could not extract text from file.")

                splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
                splits = splitter.split_documents(docs)
                
                # Assign default public access to uploaded files so the uploader can always query them
                for split in splits:
                    split.metadata["access"] = "public"

                vectorstore = Chroma.from_documents(
                    documents=splits, 
                    embedding=embeddings_model,
                    collection_name=f"temp_{file_hash}"
                )
                
                if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
                       temp_vector_stores_cache.pop(next(iter(temp_vector_stores_cache)))
                temp_vector_stores_cache[file_hash] = vectorstore
            finally:
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)

        # Temporary files (User Uploads) are treated as Admin role implicitly for that request
        # This ensures the user can query the file they just uploaded, regardless of filename.
        data = await get_rag_response(query, vectorstore, str(request.base_url), role="admin")
        return ChatResponse(query=query, **data)

    except Exception as e:
        print(f"Endpoint Error: {e}")
        raise HTTPException(500, f"Processing failed: {str(e)}")

if __name__ == "__main__":
    uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)