#semi-working agent


# # --- Force Hugging Face Transformers to be offline ---
# import os
# os.environ['HF_HUB_OFFLINE'] = '1'

# # --- API & UTILITY IMPORTS ---
# import uvicorn
# import tempfile
# import hashlib
# import shutil
# from fastapi import FastAPI, UploadFile, File, Form, HTTPException
# from fastapi.middleware.cors import CORSMiddleware
# from fastapi.responses import StreamingResponse, FileResponse
# from fastapi.staticfiles import StaticFiles
# from pydantic import BaseModel, Field # Correct Pydantic import
# from typing import List, Dict, Optional, Type
# import torch # Need torch to check for CUDA availability
# import asyncio # For running TTS in executor

# # --- TTS IMPORT ---
# from TTS.api import TTS

# # --- LANGCHAIN AND AI LIBRARIES (MODERN IMPORTS + Agent) ---
# from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader # Use PyPDFLoader
# from langchain_core.documents import Document # Needed for creating docs
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_chroma import Chroma
# from langchain_ollama import OllamaLLM
# from langchain.tools import BaseTool
# from langchain_core.callbacks import CallbackManagerForToolRun
# from langchain import hub # For agent prompts
# from langchain.agents import AgentExecutor, create_react_agent
# from langchain_core.messages import BaseMessage
# from langchain_core.vectorstores import VectorStore # Keep for type hinting

# # --- OCR SPECIFIC IMPORTS ---
# # Removed OCR imports as PyPDFLoader handles text extraction directly

# # ===================================================================
# # GLOBAL CONFIGURATION
# # ===================================================================

# CHROMA_DB_PATH = "./chroma_db_folder"
# KNOWLEDGE_BASE_PATH = r"C:\Users\vijay\OneDrive\Desktop\RAG-Base" # Use your actual path
# AUDIO_DIR = "./static_audio" # Folder to serve generated audio

# os.makedirs(AUDIO_DIR, exist_ok=True)
# os.makedirs(CHROMA_DB_PATH, exist_ok=True) # Ensure DB path exists

# # ===================================================================
# # IN-MEMORY CACHE FOR TEMPORARY VECTOR STORES
# # ===================================================================
# temp_vector_stores_cache = {}
# MAX_CACHE_SIZE = 5

# # ===================================================================
# # LOAD MODELS ON STARTUP
# # ===================================================================
# _is_cuda_available = torch.cuda.is_available()
# _device = 'cuda' if _is_cuda_available else 'cpu'
# _embeddings_model_kwargs = {'device': _device}
# _encode_kwargs = {'normalize_embeddings': False}

# print(f"CUDA Available: {_is_cuda_available}, Using device: {_device}")

# _tts_model_instance = None
# print("Loading TTS model...")
# try:
#     _tts_model_instance = TTS("tts_models/en/ljspeech/tacotron2-DDC", gpu=_is_cuda_available)
#     print("TTS model loaded successfully.")
# except Exception as e:
#     print(f"Error loading TTS model: {e}. TTS functionality will be disabled.")

# _embeddings_model = None
# print("Loading embedding model (all-MiniLM-L6-v2)...")
# try:
#     _embeddings_model = HuggingFaceEmbeddings(
#         model_name="all-MiniLM-L6-v2",
#         model_kwargs=_embeddings_model_kwargs,
#         encode_kwargs=_encode_kwargs
#     )
#     print("Embedding model loaded successfully.")
# except Exception as e:
#     print(f"CRITICAL: Failed to load embedding model: {e}")
#     exit()

# _llm = None
# print("Loading Ollama LLM (phi3)...")
# try:
#     _llm = OllamaLLM(model="phi3")
#     _llm.invoke("hello") # Test connection
#     print("Ollama LLM (phi3) connected successfully.")
# except Exception as e:
#     print(f"CRITICAL: Failed to connect to Ollama: {e}")
#     print("Please ensure the Ollama server is running and 'phi3' is installed.")
#     exit()

# # ===================================================================
# # CORE FUNCTIONS
# # ===================================================================

# async def generate_audio_file_url(text: str, base_url: str) -> Optional[str]:
#     """Generates audio, saves it, returns URL."""
#     if _tts_model_instance is None:
#         print("Skipping TTS: model not loaded.")
#         return None
#     try:
#         text_hash = hashlib.md5(text.encode()).hexdigest()
#         filename = f"{text_hash}.wav"
#         file_path = os.path.join(AUDIO_DIR, filename)
#         if not os.path.exists(file_path):
#             print(f"Generating audio: {filename}")
#             if text and text.strip():
#                  loop = asyncio.get_event_loop()
#                  await loop.run_in_executor(None, lambda: _tts_model_instance.tts_to_file(text=text, file_path=file_path))
#             else:
#                  print("Warning: Attempted TTS on empty text.")
#                  return None
#         return f"{base_url}audio/{filename}" # Use relative path for serving
#     except Exception as e:
#         print(f"Error generating audio: {e}")
#         return None


# # --- Use standard loaders, no separate OCR function needed ---
# def process_file_content(file_path: str, file_name: str) -> List[Document]:
#     """Loads text using appropriate loader (PyPDFLoader or Docx2txtLoader)."""
#     print(f"Processing content for: {file_name}")
#     try:
#         if file_name.lower().endswith('.pdf'):
#             loader = PyPDFLoader(file_path) # PyPDFLoader extracts text directly
#         elif file_name.lower().endswith('.docx'):
#             loader = Docx2txtLoader(file_path)
#         else:
#             print(f"Unsupported file type: {file_name}")
#             return []

#         docs = loader.load()
#         for doc in docs:
#             doc.metadata['source'] = file_name
#         print(f"Loaded {len(docs)} sections from {file_name}.")
#         return docs

#     except Exception as e:
#         print(f"Error loading {file_name}: {e}")
#         return [Document(page_content=f"Error loading {file_name}: {e}", metadata={'source': file_name, 'error': True})]


# # --- Function to Create/Update Vector Store ---
# def create_or_update_vector_store(docs: List[Document], vectorstore_instance: Optional[Chroma] = None, persist_path: Optional[str] = None) -> Chroma:
#     """Adds documents to an existing Chroma store or creates a new one."""
#     if not docs:
#         print("No documents provided to create/update vector store.")
#         if vectorstore_instance: return vectorstore_instance
#         else: return Chroma(persist_directory=persist_path, embedding_function=_embeddings_model) if persist_path else Chroma.from_documents([], _embeddings_model)

#     text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#     splits = text_splitter.split_documents(docs)
#     print(f"Split documents into {len(splits)} chunks.")
#     if not splits:
#         print("Warning: Splitting resulted in zero chunks.")
#         if vectorstore_instance: return vectorstore_instance
#         else: return Chroma(persist_directory=persist_path, embedding_function=_embeddings_model) if persist_path else Chroma.from_documents([], _embeddings_model)

#     if vectorstore_instance:
#         print(f"Adding {len(splits)} new chunks to existing vector store...")
#         vectorstore_instance.add_documents(splits)
#         print("Existing vector store updated.")
#         return vectorstore_instance
#     else:
#         print(f"Creating new vector store with {len(splits)} chunks...")
#         if persist_path:
#              vs = Chroma.from_documents(documents=splits, embedding=_embeddings_model, persist_directory=persist_path)
#              print(f"New persistent vector store created at {persist_path}.")
#              return vs
#         else:
#              temp_db_dir = tempfile.mkdtemp()
#              vs = Chroma.from_documents(documents=splits, embedding=_embeddings_model, persist_directory=temp_db_dir)
#              print(f"New temporary vector store created at {temp_db_dir}.")
#              # Store temp path for potential cleanup later if needed
#              vs._temp_persist_path = temp_db_dir
#              return vs


# # --- Function for loading the persistent folder store ---
# def load_or_create_folder_vector_store(folder_path: str) -> Chroma:
#     """Loads persistent store, indexes new files."""
#     print(f"Loading/Updating knowledge base: {folder_path}")
#     if not os.path.exists(folder_path):
#         os.makedirs(folder_path); print(f"Created directory: {folder_path}")

#     vectorstore = Chroma(persist_directory=CHROMA_DB_PATH, embedding_function=_embeddings_model)

#     try:
#         existing_metadatas = vectorstore.get(include=['metadatas']).get('metadatas', [])
#         indexed_files = {meta['source'] for meta in existing_metadatas if meta and 'source' in meta}
#     except Exception as e:
#         print(f"Warning: Could not get metadata from Chroma: {e}. Scanning all files.")
#         indexed_files = set()
#     print(f"Found {len(indexed_files)} previously indexed files in Chroma.")

#     current_files = {f for f in os.listdir(folder_path) if f.lower().endswith(('.pdf', '.docx'))}
#     new_files = list(current_files - indexed_files)

#     if new_files:
#         print(f"New files to process: {', '.join(new_files)}")
#         all_new_docs = []
#         for filename in new_files:
#             file_path = os.path.join(folder_path, filename)
#             try:
#                 docs = process_file_content(file_path, filename) # Use the standard loader
#                 if not docs or (len(docs) == 1 and docs[0].metadata.get('error')):
#                      print(f"Skipping {filename}: No text extracted or error during loading.")
#                      continue
#                 all_new_docs.extend(docs)
#             except Exception as e:
#                 print(f"Error processing {filename}: {e}. Skipping this file.")

#         if all_new_docs:
#             vectorstore = create_or_update_vector_store(all_new_docs, vectorstore_instance=vectorstore, persist_path=CHROMA_DB_PATH)
#             print("Knowledge base updated with new files.")
#     else:
#         print("Knowledge base is up-to-date.")
#     return vectorstore

# # ===================================================================
# # AGENT TOOL DEFINITION
# # ===================================================================
# class KnowledgeQueryInput(BaseModel):
#     query: str = Field(description="The user's question to ask the knowledge base")

# class QueryKnowledgeBaseTool(BaseTool):
#     name: str = "query_knowledge_base"
#     description: str = "Use this tool ONLY to find relevant information within the processed documents (knowledge base) that answers the user's specific question. Input MUST be the user's question."
#     args_schema: Type[BaseModel] = KnowledgeQueryInput
#     vectorstore: Chroma

#     _last_retrieved_docs: List = []

#     def _run(self, query: str, run_manager: Optional[CallbackManagerForToolRun] = None) -> str:
#         print(f"Tool '{self.name}' called with query: '{query}'")
#         try:
#             doc_count = self.vectorstore._collection.count() if hasattr(self.vectorstore, '_collection') else 0
#             if doc_count == 0:
#                  print("Vector store is empty. Tool cannot retrieve.")
#                  self._last_retrieved_docs = []
#                  return "The knowledge base is currently empty. No information found."

#             retriever = self.vectorstore.as_retriever(search_type="mmr", search_kwargs={'k': 3, 'fetch_k': 15})
#             retrieved_docs = retriever.invoke(query)
#             self._last_retrieved_docs = retrieved_docs
#             if not retrieved_docs:
#                 print("Tool found no relevant documents.")
#                 return "No relevant information found in the knowledge base."
#             else:
#                 print(f"Tool retrieved {len(retrieved_docs)} documents.")
#                 context = "\n\n---\n\n".join([doc.page_content for doc in retrieved_docs if doc.page_content])
#                 if not context.strip():
#                     print("Warning: Retrieved documents had no text content.")
#                     return "Found relevant document sections, but they contained no text."
#                 # Return context without source prefix, agent will use this
#                 return context
#         except Exception as e:
#             print(f"Error during tool execution: {e}")
#             self._last_retrieved_docs = []
#             return f"Error querying knowledge base: {e}"

#     async def _arun(self, query: str, run_manager: Optional[CallbackManagerForToolRun] = None) -> str:
#         return await asyncio.to_thread(self._run, query=query)

#     def get_last_retrieved_docs(self) -> List:
#         return self._last_retrieved_docs

# # ===================================================================
# # AGENT SETUP
# # ===================================================================
# def create_agent_executor(vectorstore: Chroma) -> tuple[AgentExecutor, QueryKnowledgeBaseTool]:
#     """Creates the agent executor with the knowledge base tool."""
#     query_tool = QueryKnowledgeBaseTool(vectorstore=vectorstore)
#     tools = [query_tool]
#     # Use the ReAct prompt suitable for Ollama/phi3
#     prompt = hub.pull("hwchase17/react")
#     agent = create_react_agent(_llm, tools, prompt) # Use the globally loaded LLM
#     agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True) # Turn verbose on for debugging
#     print("Agent Executor created.")
#     return agent_executor, query_tool

# # ===================================================================
# # FASTAPI APP INITIALIZATION & ENDPOINTS
# # ===================================================================
# app = FastAPI(
#     title="Advanced RAG Agent API",
#     description="API for the Agentic RAG Assistant (SIH)"
# )

# # --- Add CORS Middleware ---
# origins = ["*"] # Adjust for production
# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=origins,
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"],
# )

# # --- Mount Static Directory for Audio ---
# app.mount("/audio", StaticFiles(directory=AUDIO_DIR), name="audio")

# # --- Load Persistent Knowledge Base on Startup ---
# _folder_vectorstore: Optional[Chroma] = None
# try:
#     print("Initializing persistent vector store on startup...")
#     _folder_vectorstore = load_or_create_folder_vector_store(KNOWLEDGE_BASE_PATH)
#     count = _folder_vectorstore._collection.count() if _folder_vectorstore and hasattr(_folder_vectorstore, '_collection') else 0
#     print(f"Persistent vector store ready. {count} documents indexed.")
# except Exception as e:
#     print(f"CRITICAL: Failed to initialize persistent vector store on startup: {e}")
#     _folder_vectorstore = None

# # ===================================================================
# # API DATA MODELS (Pydantic)
# # ===================================================================
# class ChatRequest(BaseModel):
#     query: str

# # ===================================================================
# # API ENDPOINTS
# # ===================================================================

# @app.get("/")
# def read_root():
#     return {"status": "Advanced RAG Agent API is running"}

# @app.post("/refresh-knowledge-base")
# async def refresh_knowledge_base():
#     """Triggers a re-scan of the knowledge base folder."""
#     global _folder_vectorstore
#     try:
#         loop = asyncio.get_event_loop()
#         _folder_vectorstore = await loop.run_in_executor(None, load_or_create_folder_vector_store, KNOWLEDGE_BASE_PATH)
#         count = _folder_vectorstore._collection.count() if _folder_vectorstore and hasattr(_folder_vectorstore, '_collection') else 0
#         return {"status": "Knowledge base refreshed", "documents_indexed": count}
#     except Exception as e:
#         print(f"Error refreshing knowledge base: {e}")
#         raise HTTPException(status_code=500, detail="Error refreshing knowledge base.")

# # --- Helper for Streaming Agent Response ---
# async def stream_agent_response(agent_executor: AgentExecutor, query_tool_instance: QueryKnowledgeBaseTool, query: str, base_url: str):
#     full_response_text = ""
#     source_info = "No source found."
#     try:
#         async for chunk in agent_executor.astream({"input": query}):
#             if "actions" in chunk:
#                  for action in chunk["actions"]:
#                       action_input_safe = str(action.tool_input).replace('"', '\\"') # Basic escaping
#                       yield f'data: {{ "type": "action", "tool": "{action.tool}", "tool_input": "{action_input_safe}" }}\n\n'
#             elif "steps" in chunk:
#                  for step in chunk["steps"]:
#                       obs_str = str(step.observation).replace('"', '\\"').replace('\n', '\\n')
#                       yield f'data: {{ "type": "observation", "observation": "{obs_str}" }}\n\n'
#             elif "output" in chunk:
#                  out_str = chunk["output"].replace('"', '\\"').replace('\n', '\\n')
#                  full_response_text += chunk["output"]
#                  yield f'data: {{ "type": "stream", "text": "{out_str}" }}\n\n'

#         retrieved_docs = query_tool_instance.get_last_retrieved_docs()
#         if retrieved_docs: source_info = retrieved_docs[0].metadata.get('source', 'Unknown Document')
#         elif "No relevant information found" in full_response_text: source_info = "No information found in documents."

#         final_text_str = full_response_text.replace('"', '\\"').replace('\n', '\\n')
#         final_source_str = source_info.replace('"', '\\"').replace('\n', '\\n')
#         yield f'data: {{ "type": "final_response", "text": "{final_text_str}", "source": "{final_source_str}" }}\n\n'

#         # Generate audio AFTER final response text is confirmed
#         if full_response_text:
#              audio_url = await generate_audio_file_url(full_response_text, base_url)
#              if audio_url:
#                  yield f'data: {{ "type": "audio_url", "url": "{audio_url}" }}\n\n'

#     except Exception as e:
#         err_str = str(e).replace('"', '\\"').replace('\n', '\\n')
#         yield f'data: {{ "type": "error", "message": "An error occurred: {err_str}" }}\n\n'
#     finally:
#         yield "data: [DONE]\n\n"

# @app.post("/chat-folder")
# async def chat_with_persistent_folder(request: ChatRequest):
#     """Chat using the persistent folder knowledge base via Agent."""
#     if _folder_vectorstore is None or _folder_vectorstore._collection.count() == 0:
#         raise HTTPException(status_code=400, detail="Folder knowledge base not loaded or empty. Try refreshing first.")

#     agent_executor, query_tool_instance = create_agent_executor(_folder_vectorstore)
#     base_url = "http://127.0.0.1:8000/" # Adjust if needed
#     return StreamingResponse(stream_agent_response(agent_executor, query_tool_instance, request.query, base_url), media_type="text/event-stream")


# @app.post("/upload_and_chat_single")
# async def upload_and_chat_single(
#     query: str = Form(...),
#     file: UploadFile = File(...)
# ):
#     """Uploads a single file, processes it, and chats using it via Agent."""
#     tmp_file_path = None
#     vectorstore = None
#     file_hash = None
#     temp_db_dir_to_clean = None # Keep track of temp dir for cleanup

#     try:
#         suffix = os.path.splitext(file.filename)[1].lower()
#         if suffix not in ['.pdf', '.docx']:
#             raise HTTPException(status_code=400, detail="Unsupported file type.")

#         file_content = await file.read()
#         await file.seek(0)
#         file_hash = hashlib.md5(file_content).hexdigest()
#         print(f"Received single file upload. Hash: {file_hash}")

#         if file_hash in temp_vector_stores_cache:
#             print(f"Cache hit for {file_hash}")
#             vectorstore = temp_vector_stores_cache[file_hash]
#         else:
#             print(f"Cache miss for {file_hash}. Processing...")
#             safe_filename = "".join(c for c in file.filename if c.isalnum() or c in (' ', '.', '_')).rstrip() + suffix
#             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
#                 tmp_file.write(file_content)
#                 tmp_file_path = tmp_file.name

#             # Process and create a *temporary* persistent store
#             docs = process_file_content(tmp_file_path, safe_filename)
#             vectorstore = create_or_update_vector_store(docs, persist_path=None) # Creates temp store
#             temp_db_dir_to_clean = getattr(vectorstore, '_temp_persist_path', None) # Store path if created

#             # Cache management
#             if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
#                  oldest_key, oldest_vs = next(iter(temp_vector_stores_cache.items()))
#                  print(f"Cache full. Evicting {oldest_key}")
#                  # Cleanup evicted store's temp directory if it exists
#                  old_temp_dir = getattr(oldest_vs, '_temp_persist_path', None)
#                  if old_temp_dir and os.path.exists(old_temp_dir):
#                       shutil.rmtree(old_temp_dir)
#                       print(f"Cleaned up temp DB directory for evicted cache: {old_temp_dir}")
#                  del temp_vector_stores_cache[oldest_key]
#             temp_vector_stores_cache[file_hash] = vectorstore
#             print(f"Stored temp vector store in cache for {file_hash}")

#         # Now perform the chat using the obtained vectorstore via Agent
#         agent_executor, query_tool_instance = create_agent_executor(vectorstore)
#         base_url = "http://127.0.0.1:8000/" # Adjust if needed
#         return StreamingResponse(stream_agent_response(agent_executor, query_tool_instance, query, base_url), media_type="text/event-stream")

#     finally: # Ensure cleanup happens
#         if tmp_file_path and os.path.exists(tmp_file_path):
#             try: os.remove(tmp_file_path); print(f"Cleaned up temp file: {tmp_file_path}")
#             except Exception as cleanup_e: print(f"Error cleaning up {tmp_file_path}: {cleanup_e}")
#         # Close the file explicitly if needed, although FastAPI handles this
#         # if file and hasattr(file, 'file') and not file.file.closed:
#         #     file.file.close()

# # ===================================================================
# # RUN THE API
# # ===================================================================

# if __name__ == "__main__":
#     print("Starting Uvicorn server...")
#     uvicorn.run(
#         "__main__:app", # Use __main__ since running directly
#         host="127.0.0.1",
#         port=8000,
#         reload=True # Reload on code changes (for development)
#     )











##ocr working code


# --- Force Hugging Face Transformers to be offline ---
# import os
# os.environ['HF_HUB_OFFLINE'] = '1'

# # --- API & UTILITY IMPORTS ---
# import uvicorn
# import tempfile
# import hashlib
# import shutil
# from fastapi import FastAPI, UploadFile, File, Form, HTTPException
# from fastapi.middleware.cors import CORSMiddleware
# from fastapi.responses import JSONResponse
# from fastapi.staticfiles import StaticFiles
# from pydantic import BaseModel
# from typing import Optional, List
# import torch # Need torch to check for CUDA availability

# # --- TTS IMPORT ---
# from TTS.api import TTS

# # --- LANGCHAIN AND AI LIBRARIES (MODERN IMPORTS) ---
# from langchain_community.document_loaders import Docx2txtLoader
# from langchain_core.documents import Document # Needed for creating docs from OCR text
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_chroma import Chroma
# from langchain.prompts import PromptTemplate
# from langchain_core.output_parsers import StrOutputParser
# from langchain_ollama import OllamaLLM
# from langchain_core.vectorstores import VectorStore

# # --- OCR SPECIFIC IMPORTS ---
# import pytesseract
# from pdf2image import convert_from_path
# from PIL import Image

# # --- TESSERACT CONFIGURATION ---
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


# # ===================================================================
# # GLOBAL CONFIGURATION
# # ===================================================================

# CHROMA_DB_PATH = "./chroma_db_folder"
# KNOWLEDGE_BASE_PATH = r"C:\Users\vijay\OneDrive\Desktop\RAG-Base" # Use your actual path
# AUDIO_DIR = "./static_audio" # Folder to serve generated audio

# os.makedirs(AUDIO_DIR, exist_ok=True)

# # ===================================================================
# # IN-MEMORY CACHE FOR TEMPORARY VECTOR STORES
# # ===================================================================
# temp_vector_stores_cache = {}
# # Optional: Limit cache size to avoid excessive memory use
# MAX_CACHE_SIZE = 10 # Store up to 10 temporary vector stores

# # ===================================================================
# # LOAD MODELS ON STARTUP
# # ===================================================================

# print("Loading TTS model... (This may take a moment)")
# try:
#     # Use recommended device placement if possible
#     tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC")
#     if torch.cuda.is_available():
#         tts_model.to("cuda")
#     print("TTS model loaded successfully.")
# except Exception as e:
#     print(f"Error loading TTS model: {e}. TTS functionality will be disabled.")
#     tts_model = None

# print("Loading embedding model (all-MiniLM-L6-v2) to GPU...")
# try:
#     device = 'cuda' if torch.cuda.is_available() else 'cpu'
#     print(f"Using device: {device}")
#     model_kwargs = {'device': device}
#     encode_kwargs = {'normalize_embeddings': False} # Usually recommended for Chroma
#     embeddings_model = HuggingFaceEmbeddings(
#         model_name="all-MiniLM-L6-v2",
#         model_kwargs=model_kwargs,
#         encode_kwargs=encode_kwargs
#     )
#     print("Embedding model loaded successfully.")
# except Exception as e:
#     print(f"CRITICAL: Failed to load embedding model: {e}")
#     exit()

# print("Loading Ollama LLM (phi3)...")
# try:
#     llm = OllamaLLM(model="phi3")
#     llm.invoke("hello") # Test connection
#     print("Ollama LLM (phi3) connected successfully.")
# except Exception as e:
#     print(f"CRITICAL: Failed to connect to Ollama: {e}")
#     print("Please ensure the Ollama server is running and 'phi3' is installed.")
#     exit()

# # ===================================================================
# # REFACTORED CORE FUNCTIONS
# # ===================================================================

# def generate_audio_file_url(text: str, base_url: str) -> Optional[str]:
#     """Generates audio, saves it, returns URL."""
#     if tts_model is None:
#         print("Skipping TTS: model not loaded.")
#         return None
#     try:
#         text_hash = hashlib.md5(text.encode()).hexdigest()
#         filename = f"{text_hash}.wav"
#         file_path = os.path.join(AUDIO_DIR, filename)
#         if not os.path.exists(file_path):
#             print(f"Generating audio: {filename}")
#             # Ensure text is not empty
#             if text and text.strip():
#                  tts_model.tts_to_file(text=text, file_path=file_path)
#             else:
#                  print("Warning: Attempted TTS on empty text.")
#                  return None # Cannot generate audio for empty text
#         return f"{base_url}audio/{filename}"
#     except Exception as e:
#         print(f"Error generating audio: {e}")
#         return None

# def extract_text_with_ocr(file_path: str, file_name: str) -> List[Document]:
#     """Extracts text using OCR for PDF, or Docx2txtLoader for DOCX."""
#     extracted_docs = []
#     file_lower = file_name.lower()

#     if file_lower.endswith('.pdf'):
#         print(f"Starting OCR for PDF: {file_name}")
#         try:
#             # --- Explicit Poppler Path ---
#             images = convert_from_path(
#                 file_path,
#                 poppler_path=r"C:\Program Files\poppler-25.07.0\Library\bin" # Use your confirmed path
#             )
#             print(f"Converted {len(images)} PDF pages to images.")

#             for i, image in enumerate(images):
#                 page_num = i + 1
#                 try:
#                     # Perform OCR
#                     # Consider adding config options like --psm if needed
#                     text = pytesseract.image_to_string(image, lang='eng')
#                     print(f"  - OCR extracted text from page {page_num} (length: {len(text)})")
#                     page_doc = Document(
#                         page_content=text if text else "", # Ensure content is string
#                         metadata={'source': file_name, 'page': page_num}
#                     )
#                     extracted_docs.append(page_doc)
#                 except pytesseract.TesseractNotFoundError:
#                     print("ERROR: Tesseract executable not found or not configured.")
#                     print("Ensure Tesseract is installed and pytesseract.pytesseract.tesseract_cmd is set if needed.")
#                     raise HTTPException(status_code=500, detail="Tesseract OCR engine not found on server.")
#                 except Exception as page_e:
#                     print(f"  - Error during OCR on page {page_num}: {page_e}")
#                     extracted_docs.append(Document(page_content="", metadata={'source': file_name, 'page': page_num, 'error': str(page_e)}))

#             if not extracted_docs:
#                  print(f"Warning: OCR yielded no documents for {file_name}")

#         except Exception as e:
#             # Catch errors specifically from pdf2image/poppler if possible
#             print(f"Error converting PDF to images or during OCR process: {e}")
#             # Check if it's the specific Poppler error message
#             if "Unable to get page count" in str(e):
#                  print("Poppler error detected. Ensure Poppler's bin directory is correctly specified in poppler_path.")
#                  raise HTTPException(status_code=500, detail="Error finding or using Poppler PDF tools.")
#             else:
#                  raise HTTPException(status_code=500, detail=f"Error processing PDF for OCR: {str(e)}")


#     elif file_lower.endswith('.docx'):
#         print(f"Loading DOCX: {file_name}")
#         try:
#             loader = Docx2txtLoader(file_path)
#             docs = loader.load()
#             for doc in docs:
#                 doc.metadata['source'] = file_name # Ensure metadata
#             extracted_docs.extend(docs)
#         except Exception as e:
#             print(f"Error loading DOCX {file_name}: {e}")
#             raise HTTPException(status_code=500, detail=f"Error processing DOCX: {str(e)}")
#     else:
#         print(f"Unsupported file type: {file_name}")
#         raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_name}. Only PDF and DOCX supported.")

#     print(f"Finished extraction for {file_name}. Got {len(extracted_docs)} document sections.")
#     return extracted_docs

# def create_single_file_vector_store(file_path: str, file_name: str) -> VectorStore:
#     """Processes file (with OCR) and creates in-memory vector store."""
#     print(f"Processing vector store for: {file_name}")
#     try:
#         docs = extract_text_with_ocr(file_path, file_name)
#         if not docs:
#              print(f"No text extracted from {file_name}, creating empty vector store.")
#              return Chroma.from_documents(documents=[], embedding=embeddings_model)

#         text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#         splits = text_splitter.split_documents(docs)
#         print(f"Split {file_name} into {len(splits)} chunks.")
#         if not splits:
#              print(f"Splitting resulted in zero chunks for {file_name}, creating empty vector store.")
#              return Chroma.from_documents(documents=[], embedding=embeddings_model)

#         print(f"Creating embeddings for {len(splits)} chunks...")
#         vectorstore = Chroma.from_documents(documents=splits, embedding=embeddings_model)
#         print(f"Vector store created successfully for '{file_name}'.")
#         return vectorstore

#     except HTTPException as httpe:
#          raise httpe # Pass HTTP exceptions up
#     except Exception as e:
#         print(f"Unexpected error creating vector store for {file_name}: {e}")
#         raise HTTPException(status_code=500, detail="Internal server error during vector store creation.")


# def load_or_create_folder_vector_store(folder_path: str) -> Chroma:
#     """Loads persistent store, indexes new files (with OCR)."""
#     print(f"Loading/Updating knowledge base: {folder_path}")
#     if not os.path.exists(folder_path):
#         os.makedirs(folder_path); print(f"Created directory: {folder_path}")

#     vectorstore = Chroma(
#         persist_directory=CHROMA_DB_PATH,
#         embedding_function=embeddings_model
#     )

#     try: # Robust check for existing files
#         existing_metadatas = vectorstore.get(include=['metadatas']).get('metadatas', [])
#         indexed_files = {meta['source'] for meta in existing_metadatas if meta and 'source' in meta}
#     except Exception as e:
#         print(f"Warning: Could not get metadata from Chroma (may be empty): {e}. Will scan all files.")
#         indexed_files = set()
#     print(f"Found {len(indexed_files)} previously indexed files in Chroma.")


#     current_files = [f for f in os.listdir(folder_path) if f.lower().endswith(('.pdf', '.docx'))]
#     new_files = [f for f in current_files if f not in indexed_files]

#     if new_files:
#         print(f"New files to process: {', '.join(new_files)}")
#         all_new_splits = []
#         for filename in new_files:
#             file_path = os.path.join(folder_path, filename)
#             try:
#                 docs = extract_text_with_ocr(file_path, filename)
#                 if not docs:
#                     print(f"Skipping {filename}: No text extracted.")
#                     continue
#                 text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#                 splits = text_splitter.split_documents(docs)
#                 print(f"Processed '{filename}', {len(splits)} chunks.")
#                 all_new_splits.extend(splits)
#             except HTTPException as httpe:
#                 print(f"HTTP Error processing {filename}: {httpe.detail}. Skipping this file.")
#             except Exception as e:
#                 print(f"General Error processing {filename}: {e}. Skipping this file.")

#         if all_new_splits:
#             print(f"Adding {len(all_new_splits)} new chunks to vector store...")
#             vectorstore.add_documents(all_new_splits)
#             print("Knowledge base updated.")
#             # Consider persisting changes if needed: vectorstore.persist()
#     else:
#         print("Knowledge base is up-to-date.")
#     return vectorstore

# async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str) -> dict:
#     """Performs RAG pipeline (retrieve, generate, TTS)."""
#     # (This function remains largely the same)
#     print(f"Received query: {query}")
#     try:
#         retriever = vectorstore.as_retriever(search_type="mmr", search_kwargs={'k': 5, 'fetch_k': 20})
#         doc_count = 0
#         if isinstance(vectorstore, Chroma) and hasattr(vectorstore, '_collection'):
#              doc_count = vectorstore._collection.count()

#         if doc_count == 0:
#              print("Vector store is empty, cannot retrieve.")
#              retrieved_docs = []
#         else:
#              retrieved_docs = await retriever.ainvoke(query)

#         if not retrieved_docs:
#             print("No relevant documents found.")
#             return {"response": "Could not find relevant information.", "source": None, "audio_url": None}

#         context = "\n\n---\n\n".join([doc.page_content for doc in retrieved_docs if doc.page_content])
#         top_source_filename = retrieved_docs[0].metadata.get('source', 'Unknown') if retrieved_docs else 'Unknown'

#         # Handle case where context might be empty
#         if not context.strip():
#              print("Retrieved documents contained no usable text content.")
#              return {"response": "Found relevant document sections, but they contained no text.", "source": top_source_filename, "audio_url": None}

#         template = "Use the following context to answer the question concisely... \n\nContext: {context}\nQuestion: {question}\nAnswer:"
#         prompt = PromptTemplate.from_template(template)
#         rag_chain = prompt | llm | StrOutputParser()

#         print("Generating LLM response...")
#         response_text = await rag_chain.ainvoke({"context": context, "question": query})

#         print("Generating audio...")
#         audio_url = generate_audio_file_url(response_text, base_url)

#         print("RAG response complete.")
#         return {"response": response_text, "source": top_source_filename, "audio_url": audio_url}

#     except Exception as e:
#         print(f"Error in RAG pipeline: {e}")
#         raise e


# # ===================================================================
# # FASTAPI APP INITIALIZATION & ENDPOINTS
# # ===================================================================
# app = FastAPI(
#     title="Advanced RAG Agent API",
#     description="API for the Uni-RAG Agent"
# )

# # --- Add CORS Middleware ---
# origins = [
#     "http://localhost",
#     "http://localhost:3000",
#     "http://localhost:5173",
#     "http://localhost:8081",
#     "http://localhost:8082"
#     # Add your React app's deployed URL here later
# ]
# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=origins,
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"],
# )

# # --- Mount Static Directory ---
# app.mount("/audio", StaticFiles(directory=AUDIO_DIR), name="audio")

# # --- Load Persistent Knowledge Base ---
# try:
#     folder_vector_store = load_or_create_folder_vector_store(KNOWLEDGE_BASE_PATH)
#     doc_count = 0
#     if folder_vector_store and hasattr(folder_vector_store, '_collection'):
#          doc_count = folder_vector_store._collection.count()
#     print(f"Persistent vector store loaded. {doc_count} documents indexed.")
# except Exception as e:
#     print(f"CRITICAL: Failed to load persistent vector store: {e}")
#     folder_vector_store = None # Handle this in the endpoint

# # ===================================================================
# # API DATA MODELS (Pydantic)
# # ===================================================================
# class ChatRequest(BaseModel):
#     query: str

# class ChatResponse(BaseModel):
#     query: str
#     response: str
#     source: Optional[str] = None
#     audio_url: Optional[str] = None

# # ===================================================================
# # API ENDPOINTS
# # ===================================================================

# @app.get("/")
# def read_root():
#     return {"status": "Advanced RAG Agent API is running"}

# @app.post("/refresh-knowledge-base")
# def refresh_knowledge_base():
#     """Triggers a re-scan of the knowledge base folder."""
#     global folder_vector_store
#     try:
#         folder_vector_store = load_or_create_folder_vector_store(KNOWLEDGE_BASE_PATH)
#         count = 0
#         if folder_vector_store and hasattr(folder_vector_store, '_collection'):
#              count = folder_vector_store._collection.count()
#         return {"status": "Knowledge base refreshed", "documents_indexed": count}
#     except Exception as e:
#         print(f"Error refreshing knowledge base: {e}")
#         raise HTTPException(status_code=500, detail="Error refreshing knowledge base.")

# @app.post("/chat-folder", response_model=ChatResponse)
# async def chat_with_folder(request: ChatRequest):
#     """Chat with the persistent, pre-loaded folder of documents."""
#     if folder_vector_store is None:
#         raise HTTPException(status_code=500, detail="Store not loaded.")

#     base_url = "http://127.0.0.1:8000/" # Adjust if needed
#     try:
#         response_data = await get_rag_response(request.query, folder_vector_store, base_url)
#         return ChatResponse(query=request.query, **response_data)
#     except Exception as e:
#          print(f"Error in /chat-folder handling: {e}")
#          raise HTTPException(status_code=500, detail="Error processing request.")


# @app.post("/chat-file", response_model=ChatResponse)
# async def chat_with_single_file(
#     query: str = Form(...),
#     file: UploadFile = File(...)
# ):
#     """Upload a single file for a temporary chat session. Caches the processed vector store."""
#     tmp_file_path = None
#     vectorstore = None
#     file_hash = None

#     try:
#         suffix = os.path.splitext(file.filename)[1].lower()
#         if suffix not in ['.pdf', '.docx']:
#             raise HTTPException(status_code=400, detail="Unsupported file type.")

#         # --- Caching Step 1: Calculate File Hash ---
#         file_content = await file.read()
#         await file.seek(0) # Reset pointer
#         file_hash = hashlib.md5(file_content).hexdigest()
#         print(f"Calculated hash for {file.filename}: {file_hash}")

#         # --- Caching Step 2: Check Cache ---
#         if file_hash in temp_vector_stores_cache:
#             print(f"Cache hit for file hash: {file_hash}. Reusing vector store.")
#             vectorstore = temp_vector_stores_cache[file_hash]
#         else:
#             print(f"Cache miss for file hash: {file_hash}. Processing file...")
#             # --- Process File (Only if not in cache) ---
#             # Ensure filename is safe
#             safe_filename = "".join(c for c in file.filename if c.isalnum() or c in (' ', '.', '_')).rstrip()
#             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
#                 tmp_file.write(file_content) # Write content read earlier
#                 tmp_file_path = tmp_file.name
#             print(f"Temp file saved: {tmp_file_path}")

#             vectorstore = create_single_file_vector_store(tmp_file_path, safe_filename)

#             # --- Caching Step 3: Store in Cache ---
#             if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
#                  oldest_key = next(iter(temp_vector_stores_cache))
#                  print(f"Cache full. Evicting oldest entry: {oldest_key}")
#                  del temp_vector_stores_cache[oldest_key]
#             temp_vector_stores_cache[file_hash] = vectorstore
#             print(f"Stored vector store in cache for hash: {file_hash}")

#         # --- Use the vectorstore ---
#         base_url = "http://127.0.0.1:8000/" # Adjust if needed
#         response_data = await get_rag_response(query, vectorstore, base_url)
#         return ChatResponse(query=query, **response_data)

#     except HTTPException as httpe:
#          raise httpe # Re-raise known HTTP errors
#     except pytesseract.TesseractNotFoundError:
#          print("TesseractNotFoundError caught in endpoint.")
#          raise HTTPException(status_code=500, detail="OCR engine not found/configured.")
#     except Exception as e:
#         print(f"Unexpected error in /chat-file: {type(e).__name__}: {e}")
#         error_detail = f"Error processing file (hash: {file_hash})" if file_hash else "Error processing file."
#         raise HTTPException(status_code=500, detail=error_detail)
#     finally:
#         # Cleanup temp file ONLY if created
#         if tmp_file_path and os.path.exists(tmp_file_path):
#             try:
#                 os.remove(tmp_file_path); print(f"Cleaned up temp file: {tmp_file_path}")
#             except Exception as cleanup_e:
#                  print(f"Error cleaning up {tmp_file_path}: {cleanup_e}")
#         # Always close uploaded file object
#         if file and hasattr(file, 'file') and not file.file.closed:
#              file.file.close()


# # ===================================================================
# # RUN THE API
# # ===================================================================

# if __name__ == "__main__":
#     # Add explicit check for Tesseract command if needed, before starting server
#     try:
#         # This uses the command set at the top level
#         cmd = pytesseract.pytesseract.tesseract_cmd
#         if cmd and not os.path.exists(cmd):
#             print(f"WARNING: Configured Tesseract command not found at '{cmd}'")
#             print("OCR will likely fail. Ensure the path is correct.")
#         else:
#             print(f"Tesseract command set to: {cmd}")
#             # Optional: Verify Tesseract version
#             try:
#                  tesseract_version = pytesseract.get_tesseract_version()
#                  print(f"Successfully found Tesseract version: {tesseract_version}")
#             except Exception as tv_e:
#                  print(f"Could not get Tesseract version (check installation/PATH): {tv_e}")
#     except Exception as tess_check_e:
#         print(f"Could not verify Tesseract configuration: {tess_check_e}")
#         print("Ensure Tesseract is installed and PATH or tesseract_cmd is set.")

#     print("Starting Uvicorn server...")
#     uvicorn.run(
#         "app:app",
#         host="127.0.0.1",
#         port=8000,
#         reload=True
#     )






#Working-1:


# import os
# import shutil
# import asyncio
# import hashlib
# import tempfile
# import uvicorn
# import torch
# import glob
# from concurrent.futures import ThreadPoolExecutor
# from typing import Optional, List, Dict, Any

# # --- API & FASTAPI IMPORTS ---
# from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
# from fastapi.middleware.cors import CORSMiddleware
# from fastapi.responses import JSONResponse
# from fastapi.staticfiles import StaticFiles
# from pydantic import BaseModel

# # --- LANGCHAIN & AI IMPORTS ---
# from langchain_core.documents import Document
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_chroma import Chroma
# from langchain.prompts import PromptTemplate
# from langchain_core.output_parsers import StrOutputParser
# from langchain_ollama import OllamaLLM
# from langchain_core.vectorstores import VectorStore

# # --- OCR & VISION IMPORTS ---
# import pytesseract
# from pdf2image import convert_from_path
# from PIL import Image

# # --- DOCX IMPORT (UPDATED) ---
# try:
#     from docx import Document as DocxDocument
#     DOCX_AVAILABLE = True
# except ImportError:
#     DOCX_AVAILABLE = False
#     print("WARNING: 'python-docx' not found. .docx support disabled.")

# # --- PPTX IMPORT ---
# try:
#     from pptx import Presentation
#     PPTX_AVAILABLE = True
# except ImportError:
#     PPTX_AVAILABLE = False
#     print("WARNING: 'python-pptx' not found. .pptx support disabled.")

# # --- TTS IMPORT ---
# try:
#     from TTS.api import TTS
#     TTS_AVAILABLE = True
# except ImportError:
#     print("WARNING: 'TTS' library not found. Audio generation will be disabled.")
#     TTS_AVAILABLE = False


# # ===================================================================
# # 1. CONFIGURATION & PATH MANAGEMENT
# # ===================================================================

# class AgentConfig:
#     # 1. Tesseract Path
#     TESSERACT_CMD = shutil.which("tesseract") or r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
#     # 2. Poppler Path (For PDF -> Image)
#     POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
    
#     # 3. Storage Paths
#     CHROMA_DB_PATH = "./chroma_db_folder"
#     KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\RAG-Base") 
#     AUDIO_DIR = "./static_audio"

#     # 4. Model Settings
#     LLM_MODEL = "phi3"
#     EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
#     # 5. Supported Extensions
#     ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg'}

# # Apply Configuration
# os.environ['HF_HUB_OFFLINE'] = '1'
# os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)
# pytesseract.pytesseract.tesseract_cmd = AgentConfig.TESSERACT_CMD

# # Verify Critical Dependencies
# if not os.path.exists(AgentConfig.TESSERACT_CMD):
#     print(f"CRITICAL WARNING: Tesseract not found at {AgentConfig.TESSERACT_CMD}. OCR will fail.")
# if not os.path.exists(AgentConfig.POPPLER_PATH):
#     print(f"CRITICAL WARNING: Poppler not found at {AgentConfig.POPPLER_PATH}. PDF processing will fail.")


# # ===================================================================
# # 2. GLOBAL RESOURCES & THREAD POOL
# # ===================================================================
# thread_pool = ThreadPoolExecutor(max_workers=4)
# temp_vector_stores_cache = {}
# MAX_CACHE_SIZE = 10 

# # --- LOAD MODELS ---
# print("\n--- Loading Models ---")

# # 1. Embeddings
# print(f"Loading Embedding Model ({AgentConfig.EMBEDDING_MODEL})...")
# try:
#     device = 'cuda' if torch.cuda.is_available() else 'cpu'
#     embeddings_model = HuggingFaceEmbeddings(
#         model_name=AgentConfig.EMBEDDING_MODEL,
#         model_kwargs={'device': device},
#         encode_kwargs={'normalize_embeddings': False}
#     )
#     print("✓ Embeddings Loaded")
# except Exception as e:
#     print(f"❌ Failed to load embeddings: {e}")
#     exit(1)

# # 2. LLM
# print(f"Connecting to Ollama ({AgentConfig.LLM_MODEL})...")
# try:
#     llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
#     print("✓ Ollama Connected")
# except Exception as e:
#     print(f"❌ Failed to connect to Ollama: {e}")

# # 3. TTS
# tts_model = None
# if TTS_AVAILABLE:
#     print("Loading TTS Model...")
#     try:
#         tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC")
#         if torch.cuda.is_available():
#             tts_model.to("cuda")
#         print("✓ TTS Model Loaded")
#     except Exception as e:
#         print(f"⚠️ TTS Load Error: {e}. Audio features disabled.")
#         tts_model = None


# # ===================================================================
# # 3. CORE FUNCTIONS (Universal File Processor)
# # ===================================================================

# def _run_processing_sync(file_path: str, file_name: str) -> List[Document]:
#     """Blocking function to process PDF, DOCX, PPTX, and Images."""
#     extracted_docs = []
#     file_lower = file_name.lower()
#     ext = os.path.splitext(file_lower)[1]

#     try:
#         # --- PDF PROCESSING ---
#         if ext == '.pdf':
#             images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
#             for i, image in enumerate(images):
#                 text = pytesseract.image_to_string(image, lang='eng')
#                 if text.strip():
#                     extracted_docs.append(Document(
#                         page_content=text,
#                         metadata={'source': file_name, 'page': i + 1}
#                     ))

#         # --- DOCX PROCESSING (UPDATED to use python-docx) ---
#         elif ext == '.docx':
#             if not DOCX_AVAILABLE:
#                 raise ImportError("python-docx is not installed.")
            
#             doc = DocxDocument(file_path)
#             full_text = []
            
#             # Extract Paragraphs
#             for para in doc.paragraphs:
#                 if para.text.strip():
#                     full_text.append(para.text)
            
#             # Extract Tables
#             for table in doc.tables:
#                 for row in table.rows:
#                     row_text = [cell.text for cell in row.cells if cell.text.strip()]
#                     if row_text:
#                         full_text.append(" | ".join(row_text))

#             combined_text = "\n".join(full_text)
#             if combined_text.strip():
#                 extracted_docs.append(Document(
#                     page_content=combined_text,
#                     metadata={'source': file_name, 'page': 1}
#                 ))

#         # --- PPTX PROCESSING ---
#         elif ext == '.pptx':
#             if not PPTX_AVAILABLE:
#                 raise ImportError("python-pptx is not installed.")
#             prs = Presentation(file_path)
#             for i, slide in enumerate(prs.slides):
#                 text_runs = []
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         text_runs.append(shape.text)
#                 full_text = "\n".join(text_runs)
#                 if full_text.strip():
#                     extracted_docs.append(Document(
#                         page_content=full_text,
#                         metadata={'source': file_name, 'page': i + 1}
#                     ))

#         # --- IMAGE PROCESSING (JPG, PNG) ---
#         elif ext in ['.png', '.jpg', '.jpeg']:
#             image = Image.open(file_path)
#             text = pytesseract.image_to_string(image, lang='eng')
#             if text.strip():
#                 extracted_docs.append(Document(
#                     page_content=text,
#                     metadata={'source': file_name, 'page': 1}
#                 ))

#     except Exception as e:
#         print(f"Error processing {file_name}: {e}")
#         # Note: We return whatever we managed to extract, or empty list
    
#     return extracted_docs

# async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
#     """Async wrapper for file processing."""
#     loop = asyncio.get_running_loop()
#     return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


# def _run_tts_sync(text: str, file_path: str):
#     if tts_model:
#         tts_model.tts_to_file(text=text, file_path=file_path)

# async def generate_audio_async(text: str, base_url: str) -> Optional[str]:
#     if not tts_model or not text.strip():
#         return None
#     try:
#         text_hash = hashlib.md5(text.encode()).hexdigest()
#         filename = f"{text_hash}.wav"
#         file_path = os.path.join(AgentConfig.AUDIO_DIR, filename)
#         if not os.path.exists(file_path):
#             loop = asyncio.get_running_loop()
#             await loop.run_in_executor(thread_pool, _run_tts_sync, text, file_path)
#         return f"{base_url}audio/{filename}"
#     except Exception:
#         return None


# # ===================================================================
# # 4. RAG PIPELINE
# # ===================================================================

# async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str) -> dict:
#     try:
#         retriever = vectorstore.as_retriever(search_type="mmr", search_kwargs={'k': 4, 'fetch_k': 20})
#         retrieved_docs = await retriever.ainvoke(query)

#         if not retrieved_docs:
#             return {"response": "I couldn't find relevant information.", "source": None, "audio_url": None}

#         context_parts = []
#         sources = set()
#         for doc in retrieved_docs:
#             src = doc.metadata.get('source', 'Unknown')
#             page = doc.metadata.get('page', '?')
#             context_parts.append(f"[Source: {src}, Page: {page}]\n{doc.page_content}")
#             sources.add(src)
        
#         full_context = "\n\n".join(context_parts)
#         top_source = list(sources)[0] if sources else "Unknown"

#         template = "Use the following context to answer the question... \n\nContext: {context}\nQuestion: {question}\nAnswer:"
#         prompt = PromptTemplate.from_template(template)
#         rag_chain = prompt | llm | StrOutputParser()
        
#         response_text = await rag_chain.ainvoke({"context": full_context, "question": query})
#         audio_url = await generate_audio_async(response_text, base_url)

#         return {"response": response_text, "source": top_source, "audio_url": audio_url}

#     except Exception as e:
#         return {"response": f"Error: {str(e)}", "source": None}


# # ===================================================================
# # 5. FASTAPI APP & SMART KNOWLEDGE BASE
# # ===================================================================

# app = FastAPI(title="Uni-RAG Agent API")

# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=["*"],
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"],
# )
# app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

# folder_vector_store = None

# def init_knowledge_base():
#     """Scans RAG-Base for ALL supported files, processes them, and saves to ChromaDB."""
#     global folder_vector_store
    
#     # 1. Initialize persistent DB
#     folder_vector_store = Chroma(
#         persist_directory=AgentConfig.CHROMA_DB_PATH,
#         embedding_function=embeddings_model
#     )
    
#     if not os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
#         print(f"Warning: KB path {AgentConfig.KNOWLEDGE_BASE_PATH} does not exist.")
#         return

#     # 2. Get list of files already in DB
#     print("--- Scanning Knowledge Base for New Files ---")
#     existing_data = folder_vector_store.get()
#     existing_sources = set()
#     if existing_data and 'metadatas' in existing_data:
#         for meta in existing_data['metadatas']:
#             if meta and 'source' in meta:
#                 existing_sources.add(meta['source'])
    
#     print(f"Found {len(existing_sources)} existing files in DB.")

#     # 3. Scan directory for files (Checking all allowed extensions)
#     files_to_process = []
#     try:
#         all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
#         for f in all_files:
#             ext = os.path.splitext(f)[1].lower()
#             if ext in AgentConfig.ALLOWED_EXTENSIONS:
#                 if f not in existing_sources:
#                     files_to_process.append(f)
#     except Exception as e:
#         print(f"Error reading RAG-Base directory: {e}")

#     # 4. Process only NEW files
#     if not files_to_process:
#         print("✓ Knowledge Base is up to date.")
#     else:
#         print(f"Found {len(files_to_process)} new files to ingest: {files_to_process}")
#         new_docs = []
#         splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

#         for filename in files_to_process:
#             file_path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, filename)
#             print(f"Processing: {filename}...")
            
#             # This calls the updated logic covering DOCX, PPTX, Images, PDF
#             docs = _run_processing_sync(file_path, filename) 
            
#             if docs:
#                 split_docs = splitter.split_documents(docs)
#                 new_docs.extend(split_docs)
#             else:
#                 print(f"  Warning: No text extracted from {filename}")

#         if new_docs:
#             print(f"Adding {len(new_docs)} chunks to ChromaDB...")
#             folder_vector_store.add_documents(new_docs)
#             print("✓ Ingestion Complete!")
#         else:
#             print("No valid content found in new files.")


# # Trigger Smart Ingestion on Startup
# init_knowledge_base()


# class ChatRequest(BaseModel):
#     query: str

# class ChatResponse(BaseModel):
#     query: str
#     response: str
#     source: Optional[str] = None
#     audio_url: Optional[str] = None


# @app.post("/chat-folder", response_model=ChatResponse)
# async def chat_with_folder(request: ChatRequest, req: Request):
#     if not folder_vector_store:
#         raise HTTPException(500, "Knowledge Base not loaded.")
#     data = await get_rag_response(request.query, folder_vector_store, str(req.base_url))
#     return ChatResponse(query=request.query, **data)


# @app.post("/chat-file", response_model=ChatResponse)
# async def chat_with_single_file(
#     request: Request,
#     query: str = Form(...),
#     file: UploadFile = File(...)
# ):
#     try:
#         suffix = os.path.splitext(file.filename)[1].lower()
#         if suffix not in AgentConfig.ALLOWED_EXTENSIONS:
#             raise HTTPException(400, f"Invalid file. Allowed: {AgentConfig.ALLOWED_EXTENSIONS}")

#         content = await file.read()
#         file_hash = hashlib.md5(content).hexdigest()
        
#         if file_hash in temp_vector_stores_cache:
#             vectorstore = temp_vector_stores_cache[file_hash]
#         else:
#             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
#                 tmp.write(content)
#                 tmp_path = tmp.name
            
#             try:
#                 # This calls the updated logic covering DOCX, PPTX, Images, PDF
#                 docs = await extract_text_async(tmp_path, file.filename)
                
#                 if not docs:
#                      raise HTTPException(400, "Could not extract text from file.")

#                 splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#                 splits = splitter.split_documents(docs)
#                 vectorstore = Chroma.from_documents(documents=splits, embedding=embeddings_model)
                
#                 if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
#                       temp_vector_stores_cache.pop(next(iter(temp_vector_stores_cache)))
#                 temp_vector_stores_cache[file_hash] = vectorstore
#             finally:
#                 if os.path.exists(tmp_path):
#                     os.remove(tmp_path)

#         data = await get_rag_response(query, vectorstore, str(request.base_url))
#         return ChatResponse(query=query, **data)

#     except Exception as e:
#         print(f"Endpoint Error: {e}")
#         raise HTTPException(500, f"Processing failed: {str(e)}")

# if __name__ == "__main__":
#     uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)



#Present working:

# import os
# import shutil
# import asyncio
# import hashlib
# import tempfile
# import uvicorn
# import torch
# import glob
# import re
# from concurrent.futures import ThreadPoolExecutor
# from typing import Optional, List, Dict, Any

# # --- API & FASTAPI IMPORTS ---
# from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
# from fastapi.middleware.cors import CORSMiddleware
# from fastapi.responses import JSONResponse
# from fastapi.staticfiles import StaticFiles
# from pydantic import BaseModel

# # --- LANGCHAIN & AI IMPORTS ---
# from langchain_core.documents import Document
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_chroma import Chroma
# from langchain.prompts import PromptTemplate
# from langchain_core.output_parsers import StrOutputParser
# from langchain_ollama import OllamaLLM
# from langchain_core.vectorstores import VectorStore

# # --- OCR & VISION IMPORTS ---
# import pytesseract
# from pdf2image import convert_from_path
# from PIL import Image, ImageOps

# # --- DOCX IMPORT ---
# try:
#     from docx import Document as DocxDocument
#     DOCX_AVAILABLE = True
# except ImportError:
#     DOCX_AVAILABLE = False
#     print("WARNING: 'python-docx' not found. .docx support disabled.")

# # --- PPTX IMPORT ---
# try:
#     from pptx import Presentation
#     PPTX_AVAILABLE = True
# except ImportError:
#     PPTX_AVAILABLE = False
#     print("WARNING: 'python-pptx' not found. .pptx support disabled.")

# # --- TTS IMPORT ---
# try:
#     from TTS.api import TTS
#     TTS_AVAILABLE = True
# except ImportError:
#     print("WARNING: 'TTS' library not found. Audio generation will be disabled.")
#     TTS_AVAILABLE = False


# # ===================================================================
# # 1. CONFIGURATION & PATH MANAGEMENT
# # ===================================================================

# class AgentConfig:
#     # 1. Tesseract Path
#     TESSERACT_CMD = shutil.which("tesseract") or r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
#     # 2. Poppler Path (For PDF -> Image)
#     POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
    
#     # 3. Storage Paths
#     CHROMA_DB_PATH = "./chroma_db_folder"
#     KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\RAG-Base") 
#     AUDIO_DIR = "./static_audio"

#     # 4. Model Settings
#     LLM_MODEL = "phi3"
#     EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
#     # 5. Supported Extensions
#     ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg'}

# # Apply Configuration
# os.environ['HF_HUB_OFFLINE'] = '1'
# os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)
# pytesseract.pytesseract.tesseract_cmd = AgentConfig.TESSERACT_CMD

# # Verify Critical Dependencies
# if not os.path.exists(AgentConfig.TESSERACT_CMD):
#     print(f"CRITICAL WARNING: Tesseract not found at {AgentConfig.TESSERACT_CMD}. OCR will fail.")
# if not os.path.exists(AgentConfig.POPPLER_PATH):
#     print(f"CRITICAL WARNING: Poppler not found at {AgentConfig.POPPLER_PATH}. PDF processing will fail.")


# # ===================================================================
# # 2. GLOBAL RESOURCES & THREAD POOL
# # ===================================================================
# thread_pool = ThreadPoolExecutor(max_workers=4)
# temp_vector_stores_cache = {}
# MAX_CACHE_SIZE = 10 

# # --- LOAD MODELS ---
# print("\n--- Loading Models ---")

# # 1. Embeddings
# print(f"Loading Embedding Model ({AgentConfig.EMBEDDING_MODEL})...")
# try:
#     device = 'cuda' if torch.cuda.is_available() else 'cpu'
#     embeddings_model = HuggingFaceEmbeddings(
#         model_name=AgentConfig.EMBEDDING_MODEL,
#         model_kwargs={'device': device},
#         encode_kwargs={'normalize_embeddings': False}
#     )
#     print("✓ Embeddings Loaded")
# except Exception as e:
#     print(f"❌ Failed to load embeddings: {e}")
#     exit(1)

# # 2. LLM
# print(f"Connecting to Ollama ({AgentConfig.LLM_MODEL})...")
# try:
#     llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
#     print("✓ Ollama Connected")
# except Exception as e:
#     print(f"❌ Failed to connect to Ollama: {e}")

# # 3. TTS
# tts_model = None
# if TTS_AVAILABLE:
#     print("Loading TTS Model...")
#     try:
#         tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC")
#         if torch.cuda.is_available():
#             tts_model.to("cuda")
#         print("✓ TTS Model Loaded")
#     except Exception as e:
#         print(f"⚠️ TTS Load Error: {e}. Audio features disabled.")
#         tts_model = None


# # ===================================================================
# # 3. IMAGE PRE-PROCESSING (LEVEL 1 FIX)
# # ===================================================================

# def improve_image_quality(image: Image.Image) -> Image.Image:
#     """
#     Pre-processes an image to make it easier for Tesseract to read.
#     1. Grayscale
#     2. Upscale (if small)
#     3. Binarize (High Contrast)
#     """
#     try:
#         # 1. Convert to Grayscale
#         image = image.convert('L')
        
#         # 2. Upscale if image is too small (Tesseract struggles with small text)
#         width, height = image.size
#         if width < 1000:
#             scale_factor = 2
#             image = image.resize((width * scale_factor, height * scale_factor), Image.Resampling.LANCZOS)
        
#         # 3. Increase Contrast / Binarization (Thresholding)
#         image = image.point(lambda x: 0 if x < 140 else 255, '1')
        
#         return image
#     except Exception as e:
#         print(f"Image preprocessing warning: {e}")
#         return image  # Return original if enhancement fails

# def clean_ocr_text(text: str) -> str:
#     """Removes garbage characters often produced by OCR."""
#     text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
#     text = re.sub(r'\s+([|.,;:])\s+', r'\1 ', text)
#     return text.strip()


# # ===================================================================
# # 4. CORE PROCESSING FUNCTION (Refactored for Levels 1 & 2)
# # ===================================================================

# def _run_processing_sync(file_path: str, file_name: str) -> List[Document]:
#     """Blocking function to process PDF, DOCX, PPTX, and Images."""
#     extracted_docs = []
#     file_lower = file_name.lower()
#     ext = os.path.splitext(file_lower)[1]

#     try:
#         # --- PDF PROCESSING ---
#         if ext == '.pdf':
#             images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
#             for i, image in enumerate(images):
#                 # LEVEL 1: Improve image before reading
#                 processed_img = improve_image_quality(image)
#                 text = pytesseract.image_to_string(processed_img, lang='eng')
#                 clean = clean_ocr_text(text)
                
#                 if clean:
#                     extracted_docs.append(Document(
#                         page_content=clean,
#                         metadata={'source': file_name, 'page': i + 1}
#                     ))

#         # --- DOCX PROCESSING (LEVEL 2: Tables) ---
#         elif ext == '.docx':
#             if not DOCX_AVAILABLE:
#                 raise ImportError("python-docx is not installed.")
            
#             doc = DocxDocument(file_path)
#             full_text = []
            
#             # Extract Paragraphs
#             for para in doc.paragraphs:
#                 if para.text.strip():
#                     full_text.append(para.text)
            
#             # LEVEL 2: Extract Tables with Structure
#             for table in doc.tables:
#                 table_content = []
#                 for row in table.rows:
#                     row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
#                     if row_cells:
#                         table_content.append(" | ".join(row_cells))
                
#                 if table_content:
#                     full_text.append("\n[TABLE DATA]\n" + "\n".join(table_content) + "\n")

#             combined_text = "\n".join(full_text)
#             if combined_text.strip():
#                 extracted_docs.append(Document(
#                     page_content=combined_text,
#                     metadata={'source': file_name, 'page': 1}
#                 ))

#         # --- PPTX PROCESSING (LEVEL 2: Tables & Shapes) ---
#         elif ext == '.pptx':
#             if not PPTX_AVAILABLE:
#                 raise ImportError("python-pptx is not installed.")
#             prs = Presentation(file_path)
            
#             for i, slide in enumerate(prs.slides):
#                 slide_text = []
                
#                 for shape in slide.shapes:
#                     # Case A: Standard Text Box
#                     if hasattr(shape, "text") and shape.text.strip():
#                         slide_text.append(shape.text)
                    
#                     # Case B: Tables inside Slides (Level 2 Fix)
#                     if shape.has_table:
#                         table_rows = []
#                         for row in shape.table.rows:
#                             row_cells = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
#                             if row_cells:
#                                 table_rows.append(" | ".join(row_cells))
#                         if table_rows:
#                             slide_text.append("\n[SLIDE TABLE]\n" + "\n".join(table_rows))

#                 full_text = "\n".join(slide_text)
#                 if full_text.strip():
#                     extracted_docs.append(Document(
#                         page_content=full_text,
#                         metadata={'source': file_name, 'page': i + 1}
#                     ))

#         # --- IMAGE PROCESSING (LEVEL 1: Enhancement) ---
#         elif ext in ['.png', '.jpg', '.jpeg']:
#             image = Image.open(file_path)
#             processed_img = improve_image_quality(image)
#             custom_config = r'--oem 3 --psm 3'
#             text = pytesseract.image_to_string(processed_img, lang='eng', config=custom_config)
#             clean = clean_ocr_text(text)
#             if clean:
#                 extracted_docs.append(Document(
#                     page_content=clean,
#                     metadata={'source': file_name, 'page': 1}
#                 ))

#     except Exception as e:
#         print(f"Error processing {file_name}: {e}")
    
#     return extracted_docs

# async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
#     """Async wrapper for file processing."""
#     loop = asyncio.get_running_loop()
#     return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


# def _run_tts_sync(text: str, file_path: str):
#     if tts_model:
#         tts_model.tts_to_file(text=text, file_path=file_path)

# async def generate_audio_async(text: str, base_url: str) -> Optional[str]:
#     if not tts_model or not text.strip():
#         return None
#     try:
#         text_hash = hashlib.md5(text.encode()).hexdigest()
#         filename = f"{text_hash}.wav"
#         file_path = os.path.join(AgentConfig.AUDIO_DIR, filename)
#         if not os.path.exists(file_path):
#             loop = asyncio.get_running_loop()
#             await loop.run_in_executor(thread_pool, _run_tts_sync, text, file_path)
#         return f"{base_url}audio/{filename}"
#     except Exception:
#         return None


# # ===================================================================
# # 5. RAG PIPELINE (FIXED: Correct Source Logic)
# # ===================================================================

# async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str) -> dict:
#     try:
#         # 1. Retrieve Documents
#         retriever = vectorstore.as_retriever(search_type="mmr", search_kwargs={'k': 4, 'fetch_k': 20})
#         retrieved_docs = await retriever.ainvoke(query)

#         if not retrieved_docs:
#             return {"response": "I couldn't find relevant information.", "source": None, "audio_url": None}

#         # 2. Extract Context and Source Deterministically
#         context_parts = []
#         for doc in retrieved_docs:
#             src = doc.metadata.get('source', 'Unknown')
#             page = doc.metadata.get('page', '?')
#             context_parts.append(f"[Source: {src}, Page: {page}]\n{doc.page_content}")
        
#         full_context = "\n\n".join(context_parts)
        
#         # [FIX]: Get source from the FIRST document (the most relevant one)
#         top_source = retrieved_docs[0].metadata.get('source', 'Unknown')

#         # 3. Generate Answer
#         template = "Use the following context to answer the question... \n\nContext: {context}\nQuestion: {question}\nAnswer:"
#         prompt = PromptTemplate.from_template(template)
#         rag_chain = prompt | llm | StrOutputParser()
        
#         response_text = await rag_chain.ainvoke({"context": full_context, "question": query})
#         audio_url = await generate_audio_async(response_text, base_url)

#         return {"response": response_text, "source": top_source, "audio_url": audio_url}

#     except Exception as e:
#         print(f"RAG Error: {e}")
#         return {"response": f"Error: {str(e)}", "source": None}


# # ===================================================================
# # 6. FASTAPI APP & SMART KNOWLEDGE BASE
# # ===================================================================

# app = FastAPI(title="Uni-RAG Agent API")

# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=["*"],
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"],
# )
# app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

# folder_vector_store = None

# def init_knowledge_base():
#     """Scans RAG-Base for ALL supported files, processes them, and saves to ChromaDB."""
#     global folder_vector_store
    
#     # 1. Initialize persistent DB
#     folder_vector_store = Chroma(
#         persist_directory=AgentConfig.CHROMA_DB_PATH,
#         embedding_function=embeddings_model
#     )
    
#     if not os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
#         print(f"Warning: KB path {AgentConfig.KNOWLEDGE_BASE_PATH} does not exist.")
#         return

#     # 2. Get list of files already in DB
#     print("--- Scanning Knowledge Base for New Files ---")
#     existing_data = folder_vector_store.get()
#     existing_sources = set()
#     if existing_data and 'metadatas' in existing_data:
#         for meta in existing_data['metadatas']:
#             if meta and 'source' in meta:
#                 existing_sources.add(meta['source'])
    
#     print(f"Found {len(existing_sources)} existing files in DB.")

#     # 3. Scan directory for files
#     files_to_process = []
#     try:
#         all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
#         for f in all_files:
#             ext = os.path.splitext(f)[1].lower()
#             if ext in AgentConfig.ALLOWED_EXTENSIONS:
#                 if f not in existing_sources:
#                     files_to_process.append(f)
#     except Exception as e:
#         print(f"Error reading RAG-Base directory: {e}")

#     # 4. Process only NEW files
#     if not files_to_process:
#         print("✓ Knowledge Base is up to date.")
#     else:
#         print(f"Found {len(files_to_process)} new files to ingest: {files_to_process}")
#         new_docs = []
#         splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

#         for filename in files_to_process:
#             file_path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, filename)
#             print(f"Processing: {filename}...")
#             docs = _run_processing_sync(file_path, filename) 
            
#             if docs:
#                 split_docs = splitter.split_documents(docs)
#                 new_docs.extend(split_docs)
#             else:
#                 print(f"  Warning: No text extracted from {filename}")

#         if new_docs:
#             print(f"Adding {len(new_docs)} chunks to ChromaDB...")
#             folder_vector_store.add_documents(new_docs)
#             print("✓ Ingestion Complete!")
#         else:
#             print("No valid content found in new files.")


# # Trigger Smart Ingestion on Startup
# init_knowledge_base()


# class ChatRequest(BaseModel):
#     query: str

# class ChatResponse(BaseModel):
#     query: str
#     response: str
#     source: Optional[str] = None
#     audio_url: Optional[str] = None


# @app.post("/chat-folder", response_model=ChatResponse)
# async def chat_with_folder(request: ChatRequest, req: Request):
#     if not folder_vector_store:
#         raise HTTPException(500, "Knowledge Base not loaded.")
#     data = await get_rag_response(request.query, folder_vector_store, str(req.base_url))
#     return ChatResponse(query=request.query, **data)


# @app.post("/chat-file", response_model=ChatResponse)
# async def chat_with_single_file(
#     request: Request,
#     query: str = Form(...),
#     file: UploadFile = File(...)
# ):
#     try:
#         suffix = os.path.splitext(file.filename)[1].lower()
#         if suffix not in AgentConfig.ALLOWED_EXTENSIONS:
#             raise HTTPException(400, f"Invalid file. Allowed: {AgentConfig.ALLOWED_EXTENSIONS}")

#         content = await file.read()
#         file_hash = hashlib.md5(content).hexdigest()
        
#         if file_hash in temp_vector_stores_cache:
#             vectorstore = temp_vector_stores_cache[file_hash]
#         else:
#             with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
#                 tmp.write(content)
#                 tmp_path = tmp.name
            
#             try:
#                 docs = await extract_text_async(tmp_path, file.filename)
                
#                 if not docs:
#                      raise HTTPException(400, "Could not extract text from file.")

#                 splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#                 splits = splitter.split_documents(docs)
                
#                 # [FIX]: Unique Collection Name using file_hash
#                 # This ensures every uploaded file gets its own isolated bucket in Chroma
#                 vectorstore = Chroma.from_documents(
#                     documents=splits, 
#                     embedding=embeddings_model,
#                     collection_name=f"temp_{file_hash}"
#                 )
                
#                 if len(temp_vector_stores_cache) >= MAX_CACHE_SIZE:
#                       temp_vector_stores_cache.pop(next(iter(temp_vector_stores_cache)))
#                 temp_vector_stores_cache[file_hash] = vectorstore
#             finally:
#                 if os.path.exists(tmp_path):
#                     os.remove(tmp_path)

#         data = await get_rag_response(query, vectorstore, str(request.base_url))
#         return ChatResponse(query=query, **data)

#     except Exception as e:
#         print(f"Endpoint Error: {e}")
#         raise HTTPException(500, f"Processing failed: {str(e)}")

# if __name__ == "__main__":
#     uvicorn.run("main:app", host="127.0.0.2", port=8002, reload=True)
















#llava


# import os
# import shutil
# import asyncio
# import hashlib
# import tempfile
# import uvicorn
# import torch
# import re
# import io
# import base64
# from concurrent.futures import ThreadPoolExecutor
# from typing import Optional, List, Dict

# # --- API ---
# from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
# from fastapi.middleware.cors import CORSMiddleware
# from fastapi.staticfiles import StaticFiles
# from pydantic import BaseModel

# # --- AI & VISION ---
# from langchain_core.documents import Document
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_chroma import Chroma
# from langchain.prompts import PromptTemplate
# from langchain_core.output_parsers import StrOutputParser
# from langchain_ollama import OllamaLLM
# from langchain_core.vectorstores import VectorStore
# import ollama  # REQUIRED: pip install ollama

# # --- PDF TOOLS ---
# from pdf2image import convert_from_path
# from PIL import Image

# # --- FAST PDF PARSING ---
# try:
#     from pypdf import PdfReader
#     PYPDF_AVAILABLE = True
# except ImportError:
#     PYPDF_AVAILABLE = False

# # --- DOCX/PPTX ---
# try:
#     from docx import Document as DocxDocument
#     DOCX_AVAILABLE = True
# except ImportError:
#     DOCX_AVAILABLE = False

# try:
#     from pptx import Presentation
#     PPTX_AVAILABLE = True
# except ImportError:
#     PPTX_AVAILABLE = False

# # --- TTS ---
# try:
#     from TTS.api import TTS
#     TTS_AVAILABLE = True
# except ImportError:
#     TTS_AVAILABLE = False


# # ===================================================================
# # 1. CONFIGURATION
# # ===================================================================

# class AgentConfig:
#     # PATHS
#     # Note: We no longer strictly need Tesseract, but Poppler is still needed for PDF->Image
#     POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
#     CHROMA_DB_PATH = "./chroma_db_folder"
#     KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\RAG-Base") 
#     AUDIO_DIR = "./static_audio"

#     # MODELS
#     LLM_MODEL = "llama3"        # For thinking/chatting
#     VISION_MODEL = "llava"      # NEW: For looking at images
#     EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
#     ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg', '.webp'}

# os.environ['HF_HUB_OFFLINE'] = '1'
# os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)


# # ===================================================================
# # 2. GLOBAL RESOURCES
# # ===================================================================
# thread_pool = ThreadPoolExecutor(max_workers=4) 
# temp_vector_stores_cache = {}

# print("\n--- Loading Models ---")
# device = 'cuda' if torch.cuda.is_available() else 'cpu'

# embeddings_model = HuggingFaceEmbeddings(
#     model_name=AgentConfig.EMBEDDING_MODEL,
#     model_kwargs={'device': device},
#     encode_kwargs={'normalize_embeddings': False}
# )
# print("✓ Embeddings Loaded")

# llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
# print("✓ Ollama Connected")

# tts_model = None
# if TTS_AVAILABLE:
#     try:
#         tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC").to("cuda" if torch.cuda.is_available() else "cpu")
#         print("✓ TTS Model Loaded")
#     except:
#         print("⚠️ TTS Load Failed")


# # ===================================================================
# # 3. VISION INTELLIGENCE (The Upgrade)
# # ===================================================================

# def analyze_image_with_llava(image_bytes: bytes) -> str:
#     """
#     Sends the image directly to Ollama (Llava) to get a full description.
#     This replaces OCR.
#     """
#     try:
#         response = ollama.chat(
#             model=AgentConfig.VISION_MODEL,
#             messages=[
#                 {
#                     'role': 'user',
#                     'content': 'Analyze this image in extreme detail. Transcribe all visible text exactly. Describe any charts, tables, or visual elements.',
#                     'images': [image_bytes]
#                 }
#             ]
#         )
#         return response['message']['content']
#     except Exception as e:
#         print(f"Vision Error: {e}")
#         return ""

# def _run_processing_sync(file_path: str, file_name: str) -> List[Document]:
#     extracted_docs = []
#     ext = os.path.splitext(file_name)[1].lower()

#     try:
#         # --- PDF PROCESSING ---
#         if ext == '.pdf':
#             text_extracted = False
            
#             # 1. Try Fast Text Extraction First
#             if PYPDF_AVAILABLE:
#                 try:
#                     reader = PdfReader(file_path)
#                     for i, page in enumerate(reader.pages):
#                         text = page.extract_text()
#                         if text and len(text.strip()) > 50:
#                             extracted_docs.append(Document(
#                                 page_content=text.strip(),
#                                 metadata={'source': file_name, 'page': i + 1}
#                             ))
#                             text_extracted = True
#                 except:
#                     pass

#             # 2. Fallback to Vision Model (Llava) if PDF is an image scan
#             if not text_extracted:
#                 print(f"📷 Analyzing {file_name} with Vision AI...")
#                 # Convert PDF to images
#                 images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
                
#                 for i, image in enumerate(images):
#                     # Convert PIL image to bytes
#                     img_byte_arr = io.BytesIO()
#                     image.save(img_byte_arr, format='JPEG')
#                     img_bytes = img_byte_arr.getvalue()
                    
#                     # Ask Llava what it sees
#                     description = analyze_image_with_llava(img_bytes)
                    
#                     if description:
#                         extracted_docs.append(Document(
#                             page_content=description,
#                             metadata={'source': file_name, 'page': i + 1}
#                         ))

#         # --- IMAGE PROCESSING (Direct Vision) ---
#         elif ext in ['.png', '.jpg', '.jpeg', '.webp']:
#             print(f"👁️ Looking at {file_name}...")
#             with open(file_path, "rb") as img_file:
#                 img_bytes = img_file.read()
                
#             # Send to Llava
#             description = analyze_image_with_llava(img_bytes)
            
#             if description:
#                 extracted_docs.append(Document(
#                     page_content=description,
#                     metadata={'source': file_name, 'page': 1}
#                 ))

#         # --- DOCX ---
#         elif ext == '.docx' and DOCX_AVAILABLE:
#             doc = DocxDocument(file_path)
#             full_text = [p.text for p in doc.paragraphs if p.text.strip()]
#             for table in doc.tables:
#                 rows = [ " | ".join([c.text.strip() for c in r.cells]) for r in table.rows ]
#                 full_text.extend(rows)
#             if full_text:
#                 extracted_docs.append(Document(page_content="\n".join(full_text), metadata={'source': file_name, 'page': 1}))

#         # --- PPTX ---
#         elif ext == '.pptx' and PPTX_AVAILABLE:
#             prs = Presentation(file_path)
#             for i, slide in enumerate(prs.slides):
#                 texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
#                 if texts:
#                     extracted_docs.append(Document(page_content="\n".join(texts), metadata={'source': file_name, 'page': i+1}))

#     except Exception as e:
#         print(f"Error processing {file_name}: {e}")
    
#     return extracted_docs

# async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
#     loop = asyncio.get_running_loop()
#     return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


# # ===================================================================
# # 4. RAG & API
# # ===================================================================

# app = FastAPI(title="Uni-RAG Agent API")
# app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
# app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

# folder_vector_store = None

# @app.on_event("startup")
# async def startup_event():
#     global folder_vector_store
#     folder_vector_store = Chroma(
#         persist_directory=AgentConfig.CHROMA_DB_PATH,
#         embedding_function=embeddings_model
#     )
    
#     if os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
#         existing = set(m['source'] for m in folder_vector_store.get()['metadatas'] if m)
#         all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
#         files_to_process = [f for f in all_files 
#                             if os.path.splitext(f)[1].lower() in AgentConfig.ALLOWED_EXTENSIONS 
#                             and f not in existing]

#         if files_to_process:
#             print(f"\n🚀 Ingesting {len(files_to_process)} new files using Vision AI...")
#             splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            
#             for i, f in enumerate(files_to_process):
#                 print(f"   [{i+1}/{len(files_to_process)}] Analyzing: {f}...", end=" ", flush=True)
#                 path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, f)
#                 docs = await extract_text_async(path, f)
#                 if docs:
#                     folder_vector_store.add_documents(splitter.split_documents(docs))
#                     print("✅")
#                 else:
#                     print("⚠️ (Empty)")
#             print("\n✨ Knowledge Base Updated!\n")
#         else:
#             print("✓ KB up to date.")

# # --- ENDPOINTS ---

# class ChatRequest(BaseModel):
#     query: str

# class ChatResponse(BaseModel):
#     query: str
#     response: str
#     source: Optional[str] = None
#     audio_url: Optional[str] = None

# async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str):
#     retriever = vectorstore.as_retriever(search_type="mmr", search_kwargs={'k': 5})
#     docs = await retriever.ainvoke(query)
    
#     if not docs: return {"response": "No info found.", "source": None, "audio_url": None}

#     context = "\n\n".join([f"[Source: {d.metadata.get('source')}]\n{d.page_content}" for d in docs])
    
#     template = "Context: {context}\n\nQuestion: {question}\nAnswer:"
#     prompt = PromptTemplate.from_template(template)
#     chain = prompt | llm | StrOutputParser()
#     response = await chain.ainvoke({"context": context, "question": query})

#     audio_url = None
#     if tts_model:
#         fname = f"{hashlib.md5(response.encode()).hexdigest()}.wav"
#         fpath = os.path.join(AgentConfig.AUDIO_DIR, fname)
#         if not os.path.exists(fpath):
#             tts_model.tts_to_file(text=response, file_path=fpath)
#         audio_url = f"{base_url}audio/{fname}"

#     return {"response": response, "source": docs[0].metadata.get('source'), "audio_url": audio_url}

# @app.post("/chat-folder", response_model=ChatResponse)
# async def chat_with_folder(request: ChatRequest, req: Request):
#     if not folder_vector_store: raise HTTPException(500, "KB not loaded")
#     data = await get_rag_response(request.query, folder_vector_store, str(req.base_url))
#     return ChatResponse(query=request.query, **data)

# @app.post("/chat-file", response_model=ChatResponse)
# async def chat_with_single_file(request: Request, query: str = Form(...), file: UploadFile = File(...)):
#     suffix = os.path.splitext(file.filename)[1].lower()
#     content = await file.read()
#     fhash = hashlib.md5(content).hexdigest()

#     if fhash in temp_vector_stores_cache:
#         vstore = temp_vector_stores_cache[fhash]
#     else:
#         with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
#             tmp.write(content)
#             tmp_path = tmp.name
#         try:
#             docs = await extract_text_async(tmp_path, file.filename)
#             splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
#             vstore = Chroma.from_documents(splitter.split_documents(docs), embeddings_model, collection_name=f"temp_{fhash}")
#             temp_vector_stores_cache[fhash] = vstore
#         finally:
#             if os.path.exists(tmp_path): os.remove(tmp_path)

#     data = await get_rag_response(query, vstore, str(request.base_url))
#     return ChatResponse(query=query, **data)

# if __name__ == "__main__":
#     uvicorn.run("main:app", host="127.0.0.2", port=8002, reload=True)




#minicpm-v


import os
import shutil
import asyncio
import hashlib
import tempfile
import uvicorn
import torch
import re
import io
import base64
import cv2
import numpy as np
from concurrent.futures import ThreadPoolExecutor
from typing import Optional, List, Dict

# --- API ---
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

# --- AI & VISION ---
from langchain_core.documents import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_ollama import OllamaLLM
from langchain_core.vectorstores import VectorStore
import ollama 

# --- PDF TOOLS ---
from pdf2image import convert_from_path
from PIL import Image

# --- FAST PDF PARSING ---
try:
    from pypdf import PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

# --- DOCX/PPTX ---
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# --- TTS ---
try:
    from TTS.api import TTS
    TTS_AVAILABLE = True
except ImportError:
    TTS_AVAILABLE = False


# ===================================================================
# 1. CONFIGURATION
# ===================================================================

class AgentConfig:
    # PATHS
    POPPLER_PATH = os.getenv("POPPLER_PATH", r"C:\Program Files\poppler-25.07.0\Library\bin")
    CHROMA_DB_PATH = "./chroma_db_folder"
    KNOWLEDGE_BASE_PATH = os.getenv("KB_PATH", r"C:\Users\vijay\OneDrive\Desktop\RAG-Base") 
    AUDIO_DIR = "./static_audio"

    # MODELS
    LLM_MODEL = "llama3"        # For reasoning
    VISION_MODEL = "minicpm-v"  # UPGRADE: Best model for OCR/Text
    EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    
    ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg', '.webp'}

os.environ['HF_HUB_OFFLINE'] = '1'
os.makedirs(AgentConfig.AUDIO_DIR, exist_ok=True)


# ===================================================================
# 2. GLOBAL RESOURCES
# ===================================================================
thread_pool = ThreadPoolExecutor(max_workers=4) 
temp_vector_stores_cache = {}

print("\n--- Loading Models ---")
device = 'cuda' if torch.cuda.is_available() else 'cpu'

embeddings_model = HuggingFaceEmbeddings(
    model_name=AgentConfig.EMBEDDING_MODEL,
    model_kwargs={'device': device},
    encode_kwargs={'normalize_embeddings': False}
)
print("✓ Embeddings Loaded")

llm = OllamaLLM(model=AgentConfig.LLM_MODEL)
print("✓ Ollama Connected")

tts_model = None
if TTS_AVAILABLE:
    try:
        tts_model = TTS("tts_models/en/ljspeech/tacotron2-DDC").to("cuda" if torch.cuda.is_available() else "cpu")
        print("✓ TTS Model Loaded")
    except:
        print("⚠️ TTS Load Failed")


# ===================================================================
# 3. VISION INTELLIGENCE (Enhanced Accuracy)
# ===================================================================

def enhance_image_for_ai(image_bytes: bytes) -> bytes:
    """
    Applies 'Unsharp Masking' to deblur images before sending to AI.
    This helps significantly with screenshots and blurry photos.
    """
    try:
        # Convert bytes to OpenCV format
        nparr = np.frombuffer(image_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
        
        # 1. Denoise (Clean salt-and-pepper noise)
        img = cv2.fastNlMeansDenoisingColored(img, None, 10, 10, 7, 21)

        # 2. Sharpen (Unsharp Mask)
        gaussian = cv2.GaussianBlur(img, (0, 0), 3.0)
        unsharp_image = cv2.addWeighted(img, 1.5, gaussian, -0.5, 0)
        
        # Convert back to bytes
        _, buffer = cv2.imencode('.jpg', unsharp_image, [int(cv2.IMWRITE_JPEG_QUALITY), 95])
        return buffer.tobytes()
    except Exception as e:
        print(f"⚠️ Image enhancement failed: {e}. Using original.")
        return image_bytes

def analyze_image_with_vision_model(image_bytes: bytes) -> str:
    """
    Sends the image to MiniCPM-V for high-accuracy description.
    """
    try:
        # Step 1: Clean the image
        clean_bytes = enhance_image_for_ai(image_bytes)

        # Step 2: Query the Vision Model
        response = ollama.chat(
            model=AgentConfig.VISION_MODEL,
            messages=[
                {
                    'role': 'user',
                    # This prompt is tuned for MiniCPM-V to force OCR behavior
                    'content': 'OCR TASK: Transcribe every single piece of text in this image exactly as it appears. Also describe any charts or visual structures clearly.',
                    'images': [clean_bytes]
                }
            ]
        )
        return response['message']['content']
    except Exception as e:
        print(f"Vision Error: {e}")
        return ""

def _run_processing_sync(file_path: str, file_name: str) -> List[Document]:
    extracted_docs = []
    ext = os.path.splitext(file_name)[1].lower()

    try:
        # --- PDF PROCESSING ---
        if ext == '.pdf':
            text_extracted = False
            
            # 1. Try Fast Text Extraction First (Digital PDF)
            if PYPDF_AVAILABLE:
                try:
                    reader = PdfReader(file_path)
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text()
                        if text and len(text.strip()) > 50:
                            extracted_docs.append(Document(
                                page_content=text.strip(),
                                metadata={'source': file_name, 'page': i + 1}
                            ))
                            text_extracted = True
                except:
                    pass

            # 2. Fallback to Vision Model (MiniCPM-V) for Scanned PDFs
            if not text_extracted:
                print(f"📷 Analyzing {file_name} with MiniCPM-V...")
                images = convert_from_path(file_path, poppler_path=AgentConfig.POPPLER_PATH)
                
                for i, image in enumerate(images):
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format='JPEG')
                    img_bytes = img_byte_arr.getvalue()
                    
                    description = analyze_image_with_vision_model(img_bytes)
                    
                    if description:
                        extracted_docs.append(Document(
                            page_content=description,
                            metadata={'source': file_name, 'page': i + 1}
                        ))

        # --- IMAGE PROCESSING (Direct Vision) ---
        elif ext in ['.png', '.jpg', '.jpeg', '.webp']:
            print(f"👁️ Analyzing {file_name} with {AgentConfig.VISION_MODEL}...")
            with open(file_path, "rb") as img_file:
                img_bytes = img_file.read()
            
            description = analyze_image_with_vision_model(img_bytes)
            
            if description:
                extracted_docs.append(Document(
                    page_content=description,
                    metadata={'source': file_name, 'page': 1}
                ))

        # --- DOCX ---
        elif ext == '.docx' and DOCX_AVAILABLE:
            doc = DocxDocument(file_path)
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                rows = [ " | ".join([c.text.strip() for c in r.cells]) for r in table.rows ]
                full_text.extend(rows)
            if full_text:
                extracted_docs.append(Document(page_content="\n".join(full_text), metadata={'source': file_name, 'page': 1}))

        # --- PPTX ---
        elif ext == '.pptx' and PPTX_AVAILABLE:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                if texts:
                    extracted_docs.append(Document(page_content="\n".join(texts), metadata={'source': file_name, 'page': i+1}))

    except Exception as e:
        print(f"Error processing {file_name}: {e}")
    
    return extracted_docs

async def extract_text_async(file_path: str, file_name: str) -> List[Document]:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(thread_pool, _run_processing_sync, file_path, file_name)


# ===================================================================
# 4. RAG & API
# ===================================================================

app = FastAPI(title="Uni-RAG Agent API")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
app.mount("/audio", StaticFiles(directory=AgentConfig.AUDIO_DIR), name="audio")

folder_vector_store = None

@app.on_event("startup")
async def startup_event():
    global folder_vector_store
    folder_vector_store = Chroma(
        persist_directory=AgentConfig.CHROMA_DB_PATH,
        embedding_function=embeddings_model
    )
    
    if os.path.exists(AgentConfig.KNOWLEDGE_BASE_PATH):
        existing = set(m['source'] for m in folder_vector_store.get()['metadatas'] if m)
        all_files = os.listdir(AgentConfig.KNOWLEDGE_BASE_PATH)
        files_to_process = [f for f in all_files 
                            if os.path.splitext(f)[1].lower() in AgentConfig.ALLOWED_EXTENSIONS 
                            and f not in existing]

        if files_to_process:
            print(f"\n🚀 Ingesting {len(files_to_process)} new files using MiniCPM-V...")
            splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            
            for i, f in enumerate(files_to_process):
                print(f"   [{i+1}/{len(files_to_process)}] Analyzing: {f}...", end=" ", flush=True)
                path = os.path.join(AgentConfig.KNOWLEDGE_BASE_PATH, f)
                docs = await extract_text_async(path, f)
                if docs:
                    folder_vector_store.add_documents(splitter.split_documents(docs))
                    print("✅")
                else:
                    print("⚠️ (Empty)")
            print("\n✨ Knowledge Base Updated!\n")
        else:
            print("✓ KB up to date.")

# --- ENDPOINTS ---

class ChatRequest(BaseModel):
    query: str

class ChatResponse(BaseModel):
    query: str
    response: str
    source: Optional[str] = None
    audio_url: Optional[str] = None

async def get_rag_response(query: str, vectorstore: VectorStore, base_url: str):
    retriever = vectorstore.as_retriever(search_type="mmr", search_kwargs={'k': 5})
    docs = await retriever.ainvoke(query)
    
    if not docs: return {"response": "No info found regarding that query in the documents.", "source": None, "audio_url": None}

    context = "\n\n".join([f"[Source: {d.metadata.get('source')}]\n{d.page_content}" for d in docs])
    
    template = """
    SYSTEM: You are a helpful assistant. Use the provided context to answer the question.
    If the context contains text transcribed from an image, treat it as factual data.
    
    Context: {context}
    
    Question: {question}
    Answer:
    """
    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm | StrOutputParser()
    response = await chain.ainvoke({"context": context, "question": query})

    audio_url = None
    if tts_model:
        fname = f"{hashlib.md5(response.encode()).hexdigest()}.wav"
        fpath = os.path.join(AgentConfig.AUDIO_DIR, fname)
        if not os.path.exists(fpath):
            tts_model.tts_to_file(text=response, file_path=fpath)
        audio_url = f"{base_url}audio/{fname}"

    return {"response": response, "source": docs[0].metadata.get('source'), "audio_url": audio_url}

@app.post("/chat-folder", response_model=ChatResponse)
async def chat_with_folder(request: ChatRequest, req: Request):
    if not folder_vector_store: raise HTTPException(500, "KB not loaded")
    data = await get_rag_response(request.query, folder_vector_store, str(req.base_url))
    return ChatResponse(query=request.query, **data)

@app.post("/chat-file", response_model=ChatResponse)
async def chat_with_single_file(request: Request, query: str = Form(...), file: UploadFile = File(...)):
    suffix = os.path.splitext(file.filename)[1].lower()
    content = await file.read()
    fhash = hashlib.md5(content).hexdigest()

    if fhash in temp_vector_stores_cache:
        vstore = temp_vector_stores_cache[fhash]
    else:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            docs = await extract_text_async(tmp_path, file.filename)
            splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            vstore = Chroma.from_documents(splitter.split_documents(docs), embeddings_model, collection_name=f"temp_{fhash}")
            temp_vector_stores_cache[fhash] = vstore
        finally:
            if os.path.exists(tmp_path): os.remove(tmp_path)

    data = await get_rag_response(query, vstore, str(request.base_url))
    return ChatResponse(query=query, **data)

if __name__ == "__main__":
    uvicorn.run("main:app", host="127.0.0.2", port=8002, reload=True)